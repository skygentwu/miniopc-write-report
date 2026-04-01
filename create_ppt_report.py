#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
设备预防性维护AI智能体发展现况报告PPT生成器
浪潮集团模板版 - 精简版（35-40页）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 浪潮集团配色方案
INSPUR_BLUE = RGBColor(0, 102, 204)      # 浪潮蓝
INSPUR_ORANGE = RGBColor(255, 102, 0)    # 浪潮橙
TEXT_DARK = RGBColor(51, 51, 51)         # 正文深灰
TEXT_LIGHT = RGBColor(102, 102, 102)     # 副标题灰
WHITE = RGBColor(255, 255, 255)

# 图片路径
IMAGE_1 = "/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_15---ac3775f3-c093-4080-a4f9-dc9dddd85e6a.png"  # 设备维护演进路线图
IMAGE_2 = "/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_16---6d8e4a58-bd17-437b-9082-ec2a53499902.png"  # AI智能体技术架构图
IMAGE_3 = "/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_17---3dad5604-d9f8-4662-be94-c1d686ce45e4.png"  # 行业应用成熟度矩阵
IMAGE_4 = "/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_18---e848e927-4319-47e5-a14e-e89d7425539b.png"  # 投入期vs收益期

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9比例
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle):
    """封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加蓝色背景条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INSPUR_BLUE
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.font.name = "微软雅黑"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "微软雅黑"
    
    # 添加浪潮Logo文字
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inspur 浪潮"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    
    # 底部装饰三角形
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11), Inches(5.8), Inches(2.333), Inches(1.5))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    # 底部标语
    slogan = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.4))
    tf = slogan.text_frame
    p = tf.paragraphs[0]
    p.text = "未来，因潮澎湃  Inspur in Future"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "微软雅黑"
    
    return slide

def add_toc_slide(prs, title, items):
    """目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INSPUR_BLUE
    bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 目录项
    y_pos = 1.5
    for i, item in enumerate(items[:13]):  # 最多13个条目
        item_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(12), Inches(0.5))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_DARK
        p.font.name = "微软雅黑"
        y_pos += 0.45
    
    # 底部装饰
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.5), Inches(6), Inches(1.833), Inches(1.2))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    return slide

def add_content_slide(prs, title, content_items):
    """内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INSPUR_BLUE
    bar.line.fill.background()
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(2.5), Inches(0.4))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inspur 浪潮"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 内容
    y_pos = 1.5
    for item in content_items[:13]:  # 每页最多13个条目
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(12), Inches(0.45))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)  # 正文16pt
        p.font.color.rgb = TEXT_DARK
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
        y_pos += 0.42
    
    # 底部装饰
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.8), Inches(6.3), Inches(1.533), Inches(1))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    small_triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(6.6), Inches(0.8), Inches(0.5))
    small_triangle.fill.solid()
    small_triangle.fill.fore_color.rgb = INSPUR_BLUE
    small_triangle.line.fill.background()
    small_triangle.rotation = 180
    
    # 底部标语
    slogan = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12), Inches(0.3))
    tf = slogan.text_frame
    p = tf.paragraphs[0]
    p.text = "未来，因潮澎湃  Inspur in Future"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "微软雅黑"
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """带图片的内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INSPUR_BLUE
    bar.line.fill.background()
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(2.5), Inches(0.4))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inspur 浪潮"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 添加图片
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.6), width=Inches(11.333))
    
    # 图片说明
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.8), Inches(6.3), Inches(1.533), Inches(1))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    small_triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(6.6), Inches(0.8), Inches(0.5))
    small_triangle.fill.solid()
    small_triangle.fill.fore_color.rgb = INSPUR_BLUE
    small_triangle.line.fill.background()
    small_triangle.rotation = 180
    
    # 底部标语
    slogan = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12), Inches(0.3))
    tf = slogan.text_frame
    p = tf.paragraphs[0]
    p.text = "未来，因潮澎湃  Inspur in Future"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "微软雅黑"
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INSPUR_BLUE
    bar.line.fill.background()
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(2.5), Inches(0.4))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inspur 浪潮"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    p.font.name = "微软雅黑"
    
    # 创建表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * num_rows)).table
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.name = "微软雅黑"
        cell.fill.solid()
        cell.fill.fore_color.rgb = INSPUR_BLUE
    
    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
            cell.text_frame.paragraphs[0].font.name = "微软雅黑"
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # 底部装饰
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.8), Inches(6.3), Inches(1.533), Inches(1))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    small_triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(6.6), Inches(0.8), Inches(0.5))
    small_triangle.fill.solid()
    small_triangle.fill.fore_color.rgb = INSPUR_BLUE
    small_triangle.line.fill.background()
    small_triangle.rotation = 180
    
    # 底部标语
    slogan = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12), Inches(0.3))
    tf = slogan.text_frame
    p = tf.paragraphs[0]
    p.text = "未来，因潮澎湃  Inspur in Future"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "微软雅黑"
    
    return slide

def add_section_slide(prs, section_num, section_title):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧大色块
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = INSPUR_BLUE
    left_bar.line.fill.background()
    
    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(3), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.font.name = "微软雅黑"
    
    # 底部装饰
    triangle1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.8), Inches(6.3), Inches(1.533), Inches(1))
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 180
    
    return slide

def add_end_slide(prs):
    """结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 蓝色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = INSPUR_BLUE
    bg.line.fill.background()
    
    # 感谢语
    thanks_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(1.5))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢您的聆听"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(13.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "未来，因潮澎湃  Inspur in Future"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inspur 浪潮"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    
    return slide

def main():
    prs = create_presentation()
    
    # ===== 1. 封面页 =====
    add_title_slide(prs, 
        "设备预防性维护AI智能体发展现况报告",
        "全制造产业版 | 2026年3月")
    
    # ===== 2. 目录页 =====
    add_toc_slide(prs, "目录", [
        "执行摘要",
        "全球发展态势与竞争格局",
        "技术架构与核心能力",
        "市场规模与政策环境",
        "发展趋势与实施建议",
        "核心案例与总结"
    ])
    
    # ===== 3. 第一部分：执行摘要 =====
    add_section_slide(prs, 1, "执行摘要")
    
    add_content_slide(prs, "研究背景", [
        "• 全球设备预测性维护（PdM）市场在2024-2025年实现范式跃迁",
        "• 从『基于规则的阈值报警』向『生成式AI与多智能体协同诊断』演进",
        "• 工业Copilot与专业设备健康智能体成为全制造产业标配",
        "• 覆盖流程工业（石化、钢铁、电力）、离散制造（汽车、电子）、民生工业（食品、医药）"
    ])
    
    add_content_slide(prs, "市场趋势与规模", [
        "• 国际趋势：西门子、GE、施耐德、ABB、罗克韦尔、IBM全面推进Industrial Copilot商业化",
        "• 合作模式：自研垂直工业机理模型 + 科技巨头底层大模型算力",
        "• 国内格局：华为、阿里云、腾讯云、百度、科大讯飞、第四范式、容知日新形成生态",
        "• 全球市场150-180亿美元（2024），预计2026年突破250亿美元，CAGR 25%+",
        "• 中国市场近100亿人民币，增速30%，成为全球最大增量市场",
        "• 2025-2026年：智能体试点期；2027-2028年：多智能体协同期；2029年后：自主决策期"
    ])
    
    add_content_slide(prs, "应用行业全景与量化成效", [
        "• 流程工业（石油化工、钢铁冶金、电力能源）：适配度★★★★★",
        "  – 连续生产特征明显，停机损失巨大",
        "• 离散制造（汽车、电子、机械加工）：适配度★★★★☆",
        "  – 自动化程度高，设备关联性强",
        "• 民生工业（食品、纺织、医药）：适配度★★★☆☆",
        "  – 卫生安全要求高，柔性生产需求强",
        "• 基础设施（矿山、轨交、水务）：适配度★★★★☆",
        "  – 分布广泛，巡检难度大",
        "",
        "• 量化成效：故障监测准确率＞90%，缺陷发现提前量提升60%+",
        "• 非计划停机减少20-40%，维护成本节约15-30%，ROI回收期6-12个月"
    ])
    
    # ===== 4. 第二部分：全球发展态势 =====
    add_section_slide(prs, 2, "全球发展态势与竞争格局")
    
    add_content_slide(prs, "国际领先企业布局概览", [
        "• 市场趋势：2024-2025年实现从规则报警到AI智能体的范式跃迁",
        "• 国际工业巨头普遍采用『自研垂直工业机理模型+科技巨头算力』模式",
        "• 西门子+微软Azure：舍弗勒集团案例，停机时间↓20%，排查时间↓50%",
        "• GE Vernova+AWS：卡塔尔能源案例，意外停机↓15-20%，年节约超千万美元",
        "• 施耐德电气：巴斯夫案例，非计划停机↓近30%",
        "• ABB：SCA纸业案例，维护规划效率↑40%，寿命延长10%",
        "• 罗克韦尔：福特汽车案例，OEE提升5-8%",
        "• IBM：陶氏化学案例，巡检工作量↓40%，资产利用率↑12%"
    ])
    
    add_table_slide(prs, "国内AI服务商布局",
        ["供应商", "主打产品/技术", "标杆客户", "量化成效"],
        [
            ["华为", "盘古制造/矿山大模型", "陕煤、长安、宝武", "故障预测准确率超90%，备件↓15%"],
            ["阿里云", "ET工业大脑+通义千问", "海螺、一汽、正泰", "非计划停机↓20%，能耗↓6%"],
            ["腾讯云", "WeMake工业互联网", "宁德时代、工业富联", "维护响应时间↓30%"],
            ["百度", "文心工业版+智能体", "国家电网、恒力石化", "故障漏检率降至1%以下"],
            ["科大讯飞", "星火+工业声纹", "国能集团、美的、海尔", "噪音环境识别率85%+"],
            ["第四范式", "先知AI平台", "九江石化、宝武、燕京", "动设备预警准确率超95%"],
            ["容知日新", "SuperCare智能运维", "中石化、中广核", "诊断报告自动生成率超80%"]
        ]
    )
    
    # ===== 5. 插入行业应用成熟度矩阵图 =====
    add_image_slide(prs, "行业应用成熟度矩阵", IMAGE_3, 
        "各制造行业技术成熟度与市场规模/潜力分析")
    
    add_content_slide(prs, "流程工业：石油化工", [
        "• 设备类型：旋转机械、压缩机、泵群、反应釜、管道",
        "• 维护痛点：高温高压、易燃易爆；连续生产，非计划停机损失巨大",
        "• 标杆案例-中石化镇海炼化：阿里云+中控技术",
        "  – 大型压缩机组AI机理模型，非计划停工时间↓22%",
        "  – 年节约超千万元",
        "• 标杆案例-中石化九江石化：第四范式",
        "  – 炼化动设备群预测维护，预警准确率超95%"
    ])
    
    add_content_slide(prs, "流程工业：钢铁冶金", [
        "• 设备类型：高炉、轧机、起重设备、输送辊道",
        "• 维护痛点：高负荷连续生产；工况恶劣（高温、粉尘、重载）",
        "• 标杆案例-宝武集团：第四范式+宝信软件",
        "  – 冷轧核心产线主电机预诊断，故障率↓18%",
        "  – 备件资金占用↓10%",
        "• 标杆案例-北京科技大学轧制中心",
        "  – 『隐患+故障』双主线AI智能体，故障漏报率≤5%，准确率＞90%"
    ])
    
    add_content_slide(prs, "流程工业：电力能源", [
        "• 设备类型：风力发电机组、火电机组、变压器、电网设备",
        "• 维护痛点：设备分布广泛；新能源发电波动性大；电网稳定性要求高",
        "• 标杆案例-国家电网浙江：百度+华为",
        "  – 无人机协同巡检智能体，预测准确率96%，巡检成本↓50%",
        "• 标杆案例-国家能源集团『擎源』",
        "  – 千亿级发电大模型，41个智能体覆盖179个电站",
        "  – 半年发现缺陷2633条"
    ])
    
    add_content_slide(prs, "离散制造与基础设施", [
        "• 汽车制造：冲压机、焊接机器人、总装线、压铸机",
        "  – 比亚迪深汕基地：非计划停线时间↓30%以上",
        "  – 中国一汽红旗：AI视觉质检，识别率超99.5%",
        "  – 长安汽车：华为盘古，压铸机故障预测准确率超90%",
        "",
        "• 电子制造：宁德时代、工业富联维护响应时间↓30%",
        "",
        "• 矿山机械：陕煤集团备件库存成本↓15%",
        "",
        "• 轨道交通：中国中车高铁轴承PHM，故障率↓超15%"
    ])
    
    add_table_slide(prs, "典型业务场景与价值量化",
        ["业务场景", "解决痛点", "量化效果"],
        [
            ["设备异常早期预警", "故障发现滞后，已造成停机", "缺陷提前量↑60%+，漏报率≤5%"],
            ["故障根因智能诊断", "定位依赖专家经验，耗时长", "诊断效率↑25-40%，准确率>90%"],
            ["维护策略自动生成", "维修方案缺乏科学依据", "非计划停机减少，重复故障率下降"],
            ["备件预测与调度", "备件积压与紧急缺料并存", "库存周转率提升，紧急采购减少"],
            ["点检作业智能化", "人工点检工作量大、标准不一", "点检工作量↓60%+"],
            ["跨业务协同优化", "运维与生产、质量脱节", "OEE提升5%"],
            ["安全风险主动防控", "安全隐患发现滞后", "违章识别→告警仅15秒"]
        ]
    )
    
    # ===== 6. 第三部分：技术架构 =====
    add_section_slide(prs, 3, "技术架构与核心能力")
    
    # ===== 7. 插入设备维护演进路线图 =====
    add_image_slide(prs, "技术演进路径", IMAGE_1,
        "从定期维护到自主维护的演进历程")
    
    add_table_slide(prs, "技术演进阶段",
        ["阶段", "时间", "核心特征", "当前状态"],
        [
            ["定期维护(BM)", "2010年前", "按固定时间周期检修", "逐步淘汰"],
            ["状态监测(CM)", "2010-2015", "实时监测参数，超限报警", "广泛应用"],
            ["预测性维护(PdM)", "2015-2020", "基于数据分析预测故障", "主流应用"],
            ["智能体维护", "2020年后", "感知-诊断-决策-执行闭环", "当前跃迁期"]
        ]
    )
    
    # ===== 8. 插入AI智能体技术架构图 =====
    add_image_slide(prs, "AI智能体核心能力架构", IMAGE_2,
        "从数据采集到交互执行的分层架构")
    
    add_content_slide(prs, "AI智能体架构层次详解", [
        "• L6 交互与执行层：可视化决策看板、自然语言交互助手（Copilot）",
        "• L5 智能体协同层：设备+工艺+品质+能源+安全Agent协同",
        "• L4 分析建模层：CNN/LSTM/Transformer故障诊断、寿命预测、根因分析",
        "• L3 数据平台层：工业资源、知识经验、企业应用层集成",
        "• L2 边缘计算层：数据预处理、特征提取、毫秒级响应",
        "• L1 数据采集层：振动+温度+电流+声纹+视觉+红外热成像+油液分析"
    ])
    
    add_table_slide(prs, "关键算法与技术能力",
        ["算法类型", "适用场景", "技术特点"],
        [
            ["CNN", "振动频谱、热成像、视觉质检", "自动提取时空特征，适合图像类故障识别"],
            ["LSTM", "时间序列预测、剩余寿命预测", "捕捉长期依赖关系，适合趋势预测"],
            ["Transformer", "多传感器融合、异常检测", "自注意力机制，擅长多维数据关联"],
            ["时序大模型", "复杂工业时序数据建模", "预训练+微调，小样本场景表现优异"],
            ["知识图谱", "故障根因分析、知识检索", "将专家经验结构化，支持推理查询"],
            ["强化学习", "维护策略优化、调度决策", "在动态环境中学习最优策略"]
        ]
    )
    
    # ===== 9. 第四部分：市场规模 =====
    add_section_slide(prs, 4, "市场规模与政策环境")
    
    add_content_slide(prs, "市场规模", [
        "• 全球市场：",
        "  – 2024年市场规模：150亿-180亿美元",
        "  – 2026年预测：突破250亿美元",
        "  – 年复合增长率（CAGR）：25%以上",
        "",
        "• 中国市场：",
        "  – 2024年市场规模：近100亿元人民币",
        "  – 2025年预测增速：30%，全球最大增量市场",
        "",
        "• 增长驱动：制造业数字化转型加速、大模型技术成熟、劳动力成本上升"
    ])
    
    add_content_slide(prs, "政策支持体系", [
        "• 《推动工业领域设备更新实施方案》- 工信部",
        "  – 推进生产设备数字化转型，提升故障智能预测普及率",
        "",
        "• 国资委『人工智能+』专项行动",
        "  – 强制要求央企加快AI技术赋能主业，纳入考核",
        "",
        "• 《制造业数字化转型行动方案》- 工信部",
        "  – 中小企业数字化普及，设备上云、数据联网",
        "",
        "• 《智能制造发展规划》- 国务院",
        "  – 设备预测性维护为智能制造能力成熟度关键指标"
    ])
    
    # ===== 10. 第五部分：发展趋势 =====
    add_section_slide(prs, 5, "发展趋势与实施建议")
    
    add_table_slide(prs, "2025-2030年技术演进趋势",
        ["阶段", "时间", "核心特征", "关键能力"],
        [
            ["智能体试点期", "2025-2026", "单场景智能体应用成熟，头部企业规模化部署", "预测准确率>90%"],
            ["多智能体协同期", "2027-2028", "跨专业智能体协同成为主流", "全局最优而非单点最优"],
            ["自主决策期", "2029-2030", "AI智能体具备自主决策能力，部分场景无人运维", "自主优化，自适应学习"]
        ]
    )
    
    # ===== 11. 插入投入期vs收益期图 =====
    add_image_slide(prs, "投入期与收益期分析", IMAGE_4,
        "ROI回收期通常在6-12个月，盈亏平衡点在6个月左右")
    
    add_content_slide(prs, "分行业实施路径建议", [
        "【流程工业-石化、钢铁、电力】",
        "• 切入点：关键旋转机械（压缩机、泵、风机）",
        "• 3-6个月：高价值设备试点，建立预测模型",
        "• 6-12个月：扩展至同类设备群，建立智能运维中心",
        "",
        "【离散制造-汽车、电子】",
        "• 切入点：瓶颈设备、关键质量影响设备",
        "• 3-6个月：单条产线试点，验证OEE提升效果",
        "• 6-12个月：全厂推广，建立设备健康档案"
    ])
    
    add_content_slide(prs, "供应商选型与行动建议", [
        "【供应商选型建议】",
        "• 大型央企/国企（营收>500亿）：华为、阿里云、百度 | 安全合规、生态完整",
        "• 中型制造企业（50-500亿）：第四范式、容知日新、航天智控 | 高ROI、行业Know-how",
        "• 中小型民企（<50亿）：阿里云IoT、腾讯云WeMake | 低成本、快速部署",
        "",
        "【下一步行动】",
        "• 即时：盘点设备数据资产，明确首批试点设备",
        "• 短期：制定数据治理方案，完成供应商初步选型",
        "• 中期：启动POC试点，建立人机协同复核机制"
    ])
    
    # ===== 12. 第六部分：核心案例 =====
    add_section_slide(prs, 6, "核心案例与总结")
    
    add_table_slide(prs, "国际与国内核心案例",
        ["业主单位", "供应商", "核心场景", "量化成效"],
        [
            ["舍弗勒集团", "西门子", "产线代码生成与故障诊断", "停机时间↓20%，排查时间↓50%"],
            ["卡塔尔能源", "GE Vernova", "燃气轮机预测性维护", "意外停机↓15-20%"],
            ["中石化镇海炼化", "阿里云+中控", "压缩机预测维护", "非计划停工↓22%"],
            ["宝武集团", "第四范式+宝信", "冷轧主电机预诊断", "故障率↓18%"],
            ["国家电网浙江", "百度+华为", "无人机协同巡检", "预测准确率96%"],
            ["比亚迪深汕基地", "腾讯+自研", "机器人预测维护", "非计划停线↓30%"],
            ["中国中车", "容知日新", "高铁轴承PHM", "故障率↓超15%"]
        ]
    )
    
    add_content_slide(prs, "核心结论", [
        "• 设备预防性维护AI智能体正从技术验证走向规模化应用的关键窗口期",
        "• 覆盖从流程工业到离散制造、从原材料到终端产品的全制造产业链",
        "• 西门子Industrial Copilot、GE APM、国家能源集团『擎源』等标杆案例",
        "  表明AI智能体已在实际生产环境中展现出可量化的业务价值",
        "• 未来五年，AI智能体将从单点工具进化为系统级能力",
        "• 推动设备维护从『被动响应』走向『主动治理』",
        "• 最终实现全制造产业的『未病先治、智能自治』"
    ])
    
    add_content_slide(prs, "核心挑战与机遇", [
        "【核心挑战】",
        "• 数据质量：占投入的40%，需要大量数据治理工作",
        "• 模型泛化：跨设备、跨场景的模型迁移能力待提升",
        "• 人员接受度：传统运维人员对AI的信任和接受需要时间",
        "",
        "【战略机遇】",
        "• 率先突破这些瓶颈的企业将获得显著的竞争优势",
        "• 更低的运维成本、更高的设备可用率、更短的故障响应时间",
        "• 在数字化转型浪潮中的领先地位"
    ])
    
    # ===== 13. 结束页 =====
    add_end_slide(prs)
    
    # 保存PPT
    output_path = "/root/.openclaw/workspace/设备预防性维护AI智能体发展现况-全制造产业版.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    print(f"总页数: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    main()
