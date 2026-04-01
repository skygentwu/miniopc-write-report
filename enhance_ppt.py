from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree
import os

def create_enhanced_ppt(source_path, output_path, theme="tech"):
    """
    美化PPT：添加背景、图表、图标
    theme: tech(科技蓝), energy(能源绿), oil(石油黑金)
    """
    
    # 颜色主题配置
    themes = {
        "tech": {
            "primary": RGBColor(0, 102, 204),      # 科技蓝
            "secondary": RGBColor(0, 153, 255),    # 亮蓝
            "accent": RGBColor(255, 153, 0),       # 橙色强调
            "bg_dark": RGBColor(16, 42, 67),       # 深蓝背景
            "text_dark": RGBColor(51, 51, 51),     # 深灰文字
            "text_light": RGBColor(255, 255, 255), # 白色文字
            "gradient_start": RGBColor(0, 60, 120),
            "gradient_end": RGBColor(0, 120, 200)
        },
        "energy": {
            "primary": RGBColor(0, 128, 80),       # 能源绿
            "secondary": RGBColor(34, 139, 34),    # 森林绿
            "accent": RGBColor(255, 193, 7),       # 黄色强调
            "bg_dark": RGBColor(20, 60, 40),
            "text_dark": RGBColor(51, 51, 51),
            "text_light": RGBColor(255, 255, 255),
            "gradient_start": RGBColor(0, 80, 60),
            "gradient_end": RGBColor(0, 140, 100)
        },
        "oil": {
            "primary": RGBColor(139, 90, 43),      # 石油棕
            "secondary": RGBColor(218, 165, 32),   # 金色
            "accent": RGBColor(220, 20, 60),       # 红色强调
            "bg_dark": RGBColor(60, 40, 20),
            "text_dark": RGBColor(51, 51, 51),
            "text_light": RGBColor(255, 255, 255),
            "gradient_start": RGBColor(80, 50, 30),
            "gradient_end": RGBColor(160, 120, 60)
        }
    }
    
    colors = themes.get(theme, themes["tech"])
    
    # 创建新的演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 添加封面页
    add_cover_slide(prs, "设备健康管理智能体", "PHM Agent 完整解决方案", colors)
    
    # 添加目录页
    add_toc_slide(prs, [
        "一、PHM核心理念与价值",
        "二、系统架构设计",
        "三、关键技术解析",
        "四、落地实施方案",
        "五、监测指标体系",
        "六、典型应用场景",
        "七、组织保障与成功要素"
    ], colors)
    
    # 第一部分：核心理念
    add_section_slide(prs, "PART 01", "PHM核心理念与价值", colors)
    
    add_content_slide_with_chart(prs, "从被动维修到预测性维护", [
        "【传统模式】事后维修 vs 定期保养",
        "• 故障发生 → 停机维修 → 生产损失",
        "• 时间驱动保养 → 资源浪费",
        "",
        "【PHM模式】状态监测 → 故障诊断 → 寿命预测 → 精准维护",
        "• 实时感知设备健康状态",
        "• 提前7-30天预警故障",
        "• 在准确时间对准确部位采取准确维修"
    ], colors, chart_type="comparison")
    
    # 添加价值图表
    add_value_chart_slide(prs, "PHM核心价值指标", [
        ("非计划停机减少", "30-50%", colors["primary"]),
        ("维护成本降低", "20-30%", colors["secondary"]),
        ("设备寿命延长", "10-20%", colors["accent"]),
        ("安全事故减少", "40%", RGBColor(220, 20, 60))
    ], colors)
    
    # 第二部分：系统架构
    add_section_slide(prs, "PART 02", "系统架构设计", colors)
    
    add_architecture_slide(prs, "五层系统架构", [
        ("感知层", "传感器网络 + 边缘计算", "振动/温度/压力/电流采集"),
        ("数据层", "数据治理与存储", "时序数据库 + 数据湖"),
        ("算法层", "AI模型与知识库", "机理模型 + 深度学习"),
        ("应用层", "智能体核心功能", "监测/诊断/预测/决策"),
        ("展示层", "可视化与交互", "数字孪生 + 移动端")
    ], colors)
    
    # 添加技术栈图
    add_tech_stack_slide(prs, "关键技术栈", {
        "数据采集": ["MQTT", "Modbus", "OPC UA", "5G"],
        "存储": ["InfluxDB", "TDengine", "PostgreSQL"],
        "算法": ["LSTM", "CNN", "Transformer", "XGBoost"],
        "可视化": ["数字孪生", "3D建模", "WebGL", "移动端"]
    }, colors)
    
    # 第三部分：关键技术
    add_section_slide(prs, "PART 03", "关键技术解析", colors)
    
    add_health_index_slide(prs, "健康状态评估模型", [
        ("绿色 85-100", "健康状态", "正常维护", RGBColor(34, 139, 34)),
        ("黄色 60-84", "注意状态", "加强监测", RGBColor(255, 193, 7)),
        ("红色 <60", "危险状态", "立即检修", RGBColor(220, 20, 60))
    ], colors)
    
    add_rul_prediction_slide(prs, "RUL剩余寿命预测", [
        ("物理机理模型", "基于材料特性、应力分析、疲劳累积", "API 579标准"),
        ("数据驱动模型", "LSTM/Transformer时序预测", "大数据训练"),
        ("混合模型", "机理+数据融合，数字孪生仿真", "精度提升30%+")
    ], colors)
    
    # 第四部分：实施方案
    add_section_slide(prs, "PART 04", "落地实施方案", colors)
    
    add_timeline_slide(prs, "四阶段推进路线", [
        ("基础搭建", "1-3月", "传感器部署\n平台搭建", colors["primary"]),
        ("模型构建", "3-6月", "特征工程\n算法训练", colors["secondary"]),
        ("规模推广", "6-12月", "全面覆盖\n系统集成", colors["accent"]),
        ("持续优化", "12月+", "智能决策\n生态扩展", RGBColor(128, 128, 128))
    ], colors)
    
    # 第五部分：指标体系
    add_section_slide(prs, "PART 05", "监测指标体系", colors)
    
    add_metrics_dashboard(prs, "技术性能指标", [
        ("故障诊断准确率", "90%", "目标值", colors["primary"]),
        ("RUL预测误差", "<±15%", "精度要求", colors["secondary"]),
        ("误报率", "<5%", "控制指标", colors["accent"]),
        ("提前预警时间", ">7天", "预警能力", RGBColor(34, 139, 34)),
        ("漏报率", "<2%", "安全底线", RGBColor(220, 20, 60))
    ], colors)
    
    # 第六部分：应用场景
    add_section_slide(prs, "PART 06", "典型应用场景", colors)
    
    add_rotating_machinery_slide(prs, "旋转机械健康管理", {
        "监测参数": ["振动", "温度", "电流", "压力"],
        "诊断故障": ["轴承磨损", "转子不平衡", "叶轮磨损", "电气故障"],
        "预测维护": ["轴承RUL<30天→更换", "振动超标→48h检修"]
    }, colors)
    
    # 第七部分：保障措施
    add_section_slide(prs, "PART 07", "组织保障与成功要素", colors)
    
    add_success_factors_slide(prs, "关键成功要素", [
        "领导重视：一把手工程，资源充分保障",
        "数据质量：传感器规范安装，数据准确可靠",
        "业务融合：与EAM/ERP深度集成",
        "持续运营：专人维护模型，定期评估优化"
    ], colors)
    
    # 添加总结页
    add_summary_slide(prs, "实施路径图", [
        ("短期目标", "6个月", "诊断准确率>85%\n停机减少20%", colors["primary"]),
        ("中期目标", "12个月", "诊断准确率>90%\n停机减少40%", colors["secondary"]),
        ("长期目标", "24个月", "AI自主决策\n行业标杆水平", colors["accent"])
    ], colors)
    
    # 添加结束页
    add_end_slide(prs, "让设备会说话、能预测、自维护", "设备健康管理智能体", colors)
    
    # 保存
    prs.save(output_path)
    print(f"✅ 美化版PPT已生成: {output_path}")
    print(f"📊 共 {len(prs.slides)} 页")
    return output_path

# ===== 幻灯片创建函数 =====

def add_cover_slide(prs, title, subtitle, colors):
    """封面页 - 带渐变背景"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加渐变背景矩形
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["bg_dark"]
    bg_shape.line.fill.background()
    
    # 装饰线条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors["accent"]
    line.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = colors["text_light"]
    
    # 副标题
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = colors["secondary"]
        p2.space_before = Pt(20)

def add_toc_slide(prs, items, colors):
    """目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目录 CONTENTS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]
    
    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.5), Inches(0.08), Inches(5.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["primary"]
    bar.line.fill.background()
    
    # 目录项
    y_pos = 1.5
    for i, item in enumerate(items, 1):
        # 编号
        num_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(0.8), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{i}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos), Inches(10), Inches(0.6))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = colors["text_dark"]
        
        y_pos += 0.7

def add_section_slide(prs, part_num, title, colors):
    """章节过渡页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["bg_dark"]
    bg_shape.line.fill.background()
    
    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = part_num
    p.font.size = Pt(28)
    p.font.color.rgb = colors["secondary"]
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors["text_light"]
    
    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors["accent"]
    line.line.fill.background()

def add_content_slide_with_chart(prs, title, bullets, colors, chart_type="default"):
    """内容页 - 带图表占位"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = colors["primary"]
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        if bullet.startswith("【"):
            p.font.bold = True
            p.font.color.rgb = colors["primary"]
    
    # 右侧图表占位框
    chart_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.3), Inches(4.8), Inches(5.5))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    chart_box.line.color.rgb = colors["primary"]
    chart_box.line.width = Pt(2)
    
    # 图表说明
    tf = chart_box.text_frame
    tf.paragraphs[0].text = "[对比示意图]" if chart_type == "comparison" else "[数据图表]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_value_chart_slide(prs, title, items, colors):
    """价值指标图表页 - 横向条形图"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]
    
    # 创建条形图
    y_pos = 1.5
    for label, value, bar_color in items:
        # 标签
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(3), Inches(0.5))
        tf = label_box.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(16)
        
        # 数值
        value_box = slide.shapes.add_textbox(Inches(3.8), Inches(y_pos), Inches(1.5), Inches(0.5))
        tf = value_box.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = bar_color
        
        # 条形
        bar_width = 2 + float(value.replace('%', '').replace('-', '')) / 20  # 根据数值计算长度
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(y_pos + 0.05), Inches(bar_width), Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()
        
        y_pos += 0.9
    
    # 添加说明
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(1))
    tf = note_box.text_frame
    tf.paragraphs[0].text = "* 数据来源于行业标杆案例统计分析"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

def add_architecture_slide(prs, title, layers, colors):
    """架构图 - 垂直分层"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 垂直分层
    y_pos = 1.1
    layer_height = 1.0
    width = 12
    
    for i, (layer_name, subtitle, desc) in enumerate(layers):
        # 层背景
        layer_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(width), Inches(layer_height))
        layer_bg.fill.solid()
        # 渐变色效果
        color_intensity = 1 - (i * 0.15)
        r = int(colors["primary"].rgb[0] * color_intensity + 240 * (1 - color_intensity))
        g = int(colors["primary"].rgb[1] * color_intensity + 240 * (1 - color_intensity))
        b = int(colors["primary"].rgb[2] * color_intensity + 240 * (1 - color_intensity))
        layer_bg.fill.fore_color.rgb = RGBColor(r, g, b)
        layer_bg.line.color.rgb = colors["primary"]
        
        # 层名称
        name_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(2.5), Inches(0.4))
        tf = name_box.text_frame
        tf.paragraphs[0].text = layer_name
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["text_dark"]
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.5), Inches(2.5), Inches(0.4))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos + 0.2), Inches(8.5), Inches(0.6))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = colors["text_dark"]
        
        y_pos += layer_height + 0.15
        
        # 连接箭头（除了最后一层）
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6), Inches(y_pos - 0.1), Inches(0.5), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = colors["accent"]

def add_tech_stack_slide(prs, title, tech_dict, colors):
    """技术栈图"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 四列布局
    categories = list(tech_dict.keys())
    col_width = 3
    x_start = 0.5
    
    for i, category in enumerate(categories):
        x_pos = x_start + i * col_width
        
        # 类别标题
        cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.2), Inches(2.7), Inches(0.5))
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = colors["primary"]
        cat_box.line.fill.background()
        
        tf = cat_box.text_frame
        tf.paragraphs[0].text = category
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["text_light"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 技术项
        y_pos = 1.9
        for tech in tech_dict[category]:
            tech_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.1), Inches(y_pos), Inches(2.5), Inches(0.4))
            tech_box.fill.solid()
            tech_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
            tech_box.line.color.rgb = colors["primary"]
            
            tf = tech_box.text_frame
            tf.paragraphs[0].text = tech
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            y_pos += 0.55

def add_health_index_slide(prs, title, levels, colors):
    """健康指数仪表盘"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 三个状态卡片
    card_width = 3.8
    gap = 0.5
    x_start = 0.5
    
    for i, (label, status, action, color) in enumerate(levels):
        x_pos = x_start + i * (card_width + gap)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.5), Inches(card_width), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # 颜色条
        color_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.5), Inches(card_width), Inches(0.6))
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = color
        color_bar.line.fill.background()
        
        # 标签
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.3), Inches(card_width), Inches(0.6))
        tf = label_box.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 状态
        status_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.2), Inches(card_width), Inches(0.5))
        tf = status_box.text_frame
        tf.paragraphs[0].text = status
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 建议动作
        action_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4), Inches(card_width - 0.4), Inches(1.5))
        tf = action_box.text_frame
        tf.paragraphs[0].text = action
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_rul_prediction_slide(prs, title, methods, colors):
    """RUL预测方法对比"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 三种方法横向排列
    card_width = 4
    gap = 0.2
    x_start = 0.5
    
    for i, (method, desc, feature) in enumerate(methods):
        x_pos = x_start + i * (card_width + gap)
        
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.3), Inches(card_width), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = colors["primary"]
        
        # 方法名
        name_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.5), Inches(card_width), Inches(0.6))
        tf = name_box.text_frame
        tf.paragraphs[0].text = method
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["primary"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.3), Inches(card_width - 0.4), Inches(2))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].word_wrap = True
        
        # 特点标签
        feature_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.5), Inches(5), Inches(3), Inches(0.6))
        feature_box.fill.solid()
        feature_box.fill.fore_color.rgb = colors["secondary"]
        feature_box.line.fill.background()
        
        tf = feature_box.text_frame
        tf.paragraphs[0].text = feature
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["text_light"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_timeline_slide(prs, title, phases, colors):
    """时间线"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 水平时间线
    timeline_y = 3
    total_width = 12
    start_x = 0.5
    
    # 主线
    main_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(timeline_y), Inches(total_width), Inches(0.05))
    main_line.fill.solid()
    main_line.fill.fore_color.rgb = colors["primary"]
    main_line.line.fill.background()
    
    # 阶段节点
    node_count = len(phases)
    spacing = total_width / (node_count - 1) if node_count > 1 else total_width
    
    for i, (phase, time, desc, color) in enumerate(phases):
        x_pos = start_x + i * spacing
        
        # 节点圆
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos - 0.2), Inches(timeline_y - 0.2), Inches(0.45), Inches(0.45))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        
        # 阶段名（上方）
        phase_box = slide.shapes.add_textbox(Inches(x_pos - 1), Inches(timeline_y - 1.2), Inches(2), Inches(0.5))
        tf = phase_box.text_frame
        tf.paragraphs[0].text = phase
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 时间（下方）
        time_box = slide.shapes.add_textbox(Inches(x_pos - 0.8), Inches(timeline_y + 0.4), Inches(1.6), Inches(0.4))
        tf = time_box.text_frame
        tf.paragraphs[0].text = time
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = colors["text_dark"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 描述（下方）
        desc_box = slide.shapes.add_textbox(Inches(x_pos - 1.2), Inches(timeline_y + 0.9), Inches(2.4), Inches(1.5))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_metrics_dashboard(prs, title, metrics, colors):
    """指标仪表盘"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 2x3 网格布局
    card_width = 4
    card_height = 2.8
    gap_x = 0.3
    gap_y = 0.4
    
    positions = [
        (0.5, 1.2), (4.8, 1.2), (9.1, 1.2),
        (0.5, 4.4), (4.8, 4.4)
    ]
    
    for i, (metric_name, value, unit, color) in enumerate(metrics):
        if i >= len(positions):
            break
        x, y = positions[i]
        
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 指标名
        name_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(card_width), Inches(0.5))
        tf = name_box.text_frame
        tf.paragraphs[0].text = metric_name
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = colors["text_dark"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 数值
        value_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.8), Inches(card_width), Inches(0.8))
        tf = value_box.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 单位说明
        unit_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.7), Inches(card_width), Inches(0.5))
        tf = unit_box.text_frame
        tf.paragraphs[0].text = unit
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_rotating_machinery_slide(prs, title, content, colors):
    """旋转机械场景"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 三列布局
    sections = [
        ("监测参数", content["监测参数"], colors["primary"]),
        ("诊断故障", content["诊断故障"], colors["secondary"]),
        ("预测维护", content["预测维护"], colors["accent"])
    ]
    
    col_width = 4
    gap = 0.2
    x_start = 0.5
    
    for i, (section_title, items, color) in enumerate(sections):
        x_pos = x_start + i * (col_width + gap)
        
        # 标题栏
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(1.2), Inches(col_width), Inches(0.6))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        
        tf = header.text_frame
        tf.paragraphs[0].text = section_title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["text_light"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 内容项
        y_pos = 2.0
        for item in items:
            item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.1), Inches(y_pos), Inches(col_width - 0.2), Inches(0.7))
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
            item_box.line.color.rgb = color
            
            tf = item_box.text_frame
            tf.paragraphs[0].text = item
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].word_wrap = True
            
            y_pos += 0.9

def add_success_factors_slide(prs, title, factors, colors):
    """成功要素"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["primary"]
    
    # 四个象限
    positions = [(0.5, 1.3), (6.8, 1.3), (0.5, 4.3), (6.8, 4.3)]
    icons = ["👨‍💼", "📊", "🔗", "🔄"]
    
    for i, (factor, (x, y)) in enumerate(zip(factors, positions)):
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = colors["primary"]
        
        # 图标
        icon_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        tf.paragraphs[0].text = icons[i]
        tf.paragraphs[0].font.size = Pt(36)
        
        # 文字
        text_box = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.4), Inches(4.5), Inches(1.8))
        tf = text_box.text_frame
        tf.paragraphs[0].text = factor
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].word_wrap = True

def add_summary_slide(prs, title, goals, colors):
    """总结页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["bg_dark"]
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["text_light"]
    
    # 三个目标卡片
    card_width = 4
    gap = 0.3
    x_start = 0.5
    
    for i, (goal_title, time, desc, color) in enumerate(goals):
        x_pos = x_start + i * (card_width + gap)
        
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2), Inches(card_width), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # 时间标签
        time_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.3), Inches(2.3), Inches(1.5), Inches(0.5))
        time_label.fill.solid()
        time_label.fill.fore_color.rgb = color
        time_label.line.fill.background()
        
        tf = time_label.text_frame
        tf.paragraphs[0].text = time
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = colors["text_light"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 标题
        goal_title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.2), Inches(card_width), Inches(0.6))
        tf = goal_title_box.text_frame
        tf.paragraphs[0].text = goal_title
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4), Inches(card_width - 0.4), Inches(2.5))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_end_slide(prs, slogan, title, colors):
    """结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors["bg_dark"]
    bg.line.fill.background()
    
    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3), Inches(5.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors["accent"]
    line.line.fill.background()
    
    # 标语
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = slogan_box.text_frame
    tf.paragraphs[0].text = slogan
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["text_light"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = colors["secondary"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 底部装饰
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.8), Inches(5.333), Inches(0.05))
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = colors["accent"]
    bottom_line.line.fill.background()

# ===== 执行美化 =====
if __name__ == "__main__":
    # 美化设备健康管理智能体PPT
    output_path = "/root/.openclaw/workspace/设备健康管理智能体-美化版.pptx"
    create_enhanced_ppt(None, output_path, theme="tech")
    print(f"✅ 美化完成: {output_path}")
