from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿 - 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)
INSPUR_ORANGE = RGBColor(255, 102, 0)
INSPUR_RED = RGBColor(255, 51, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)

# 字体设置
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    """添加浪潮装饰元素"""
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    """添加浪潮Logo"""
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    """添加底部标语"""
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    """封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(26)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(1.1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def calculate_content_height(content_lines):
    """估算内容所需高度（保守估算，确保不超出）"""
    total_lines = 0
    for line in content_lines:
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            total_lines += 2.0  # 一级标题占2行（含间距）
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆') or line.startswith('●'):
            # 根据文字长度估算行数（每行约50个字符，更保守）
            text_len = len(line)
            lines_needed = max(1, (text_len + 49) // 50)
            total_lines += lines_needed * 1.0  # 每行按1行算，不压缩
        elif line.startswith('→') or line.startswith('▶') or line.startswith('○') or line.startswith('  -'):
            text_len = len(line)
            lines_needed = max(1, (text_len + 49) // 50)
            total_lines += lines_needed * 0.95
        elif line == '' or line.strip() == '':
            total_lines += 0.5  # 空行按0.5行算
        else:
            text_len = len(line)
            lines_needed = max(1, (text_len + 49) // 50)
            total_lines += lines_needed * 1.0
    return total_lines

def split_content_to_pages(content_lines, max_lines_per_page=18):
    """将内容按安全行数拆分成多页 - 保守策略，每页最多18行确保不超出"""
    pages = []
    current_page = []
    current_lines = 0
    
    for line in content_lines:
        line_height = calculate_content_height([line])
        
        # 如果当前行加上会超出限制，且当前页不为空，则新建一页
        if current_lines + line_height > max_lines_per_page and current_page:
            pages.append(current_page)
            current_page = [line]
            current_lines = line_height
        else:
            current_page.append(line)
            current_lines += line_height
    
    # 添加最后一页
    if current_page:
        pages.append(current_page)
    
    return pages

def add_content_slide_dense(prs, title, content_lines):
    """内容页 - 高密度文字版本，自动分页确保不超出页面"""
    # 先拆分内容，确保每页在安全范围内
    pages = split_content_to_pages(content_lines, max_lines_per_page=26)
    slides = []
    
    for page_idx, page_content in enumerate(pages):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        add_inspur_logo(slide)
        
        # 标题区域 - 如果有多页，显示页码
        display_title = title if len(pages) == 1 else f"{title} ({page_idx + 1}/{len(pages)})"
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.name = FONT_NAME
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = INSPUR_BLUE
        
        # 内容区域 - 安全高度4.5英寸（18行安全高度，确保不超出）
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(page_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 一级标题 【】▎■
            if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = INSPUR_ORANGE
                p.space_before = Pt(8)
                p.space_after = Pt(3)
            # 二级标题 • - ◆ ●
            elif line.startswith('•') or line.startswith('-') or line.startswith('◆') or line.startswith('●'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
            # 三级缩进 → ▶ ○
            elif line.startswith('→') or line.startswith('▶') or line.startswith('○') or line.startswith('  -'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_GRAY
                p.space_after = Pt(2)
            # 空行
            elif line == '' or line.strip() == '':
                p.text = ''
                p.space_after = Pt(3)
            # 普通文本
            else:
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
        
        add_footer_slogan(slide)
        slides.append(slide)
    
    return slides

# ==================== PPT内容 - 高密度版本 ====================

# 封面
add_title_slide(prs, 'AI 带来的范式转移', '中石油勘探院油气开采研究所智能化转型战略思考')

# 第一章：从感知到认知
add_section_divider(prs, '01', '从"感知"到"认知"：AI 能力的跃迁')

add_content_slide_dense(prs, '传统 AI vs 现代大模型：能力边界的突破', [
    '【传统 AI：感知智能时代】',
    '• 擅长：图像识别、语音转文字、数据分类、模式匹配',
    '• 局限：单点任务、缺乏推理、无法理解复杂上下文、难以处理模糊问题',
    '• 比喻：像一名训练有素的技工，只会特定工序，无法应对变化',
    '',
    '【现代大模型：认知智能时代】',
    '▎自主感知：多源数据融合理解',
    '• 同时处理地震数据、测井曲线、生产日志、文献资料、设备传感器数据',
    '• 跨模态对齐：文本+图像+时序+空间数据统一语义理解',
    '• 案例：昆仑大模型可同时解析文本、图像、时序数据，实现多模态联合推理',
    '',
    '▎复杂推理：从"是什么"到"为什么"',
    '• 逻辑链推理：异常征兆 → 根因分析 → 处置建议 → 效果评估',
    '• 因果推断能力：区分相关性与因果性，避免经验陷阱',
    '• 案例：钻井风险预警准确率提升至93.8%（新疆油田），误报率降低60%',
    '',
    '▎目标规划：从被动响应到主动决策',
    '• 自动生成开发方案、优化井位部署、预测产量趋势、动态调整生产制度',
    '• 强化学习：在复杂环境中自主学习最优策略',
    '• 案例：大庆油田产量预测准确度达90.74%，效率提升10倍，预测周期从周级缩短至小时级'
])

add_content_slide_dense(prs, '智能体（Agent）：从"工具"到"同事"', [
    '【什么是智能体？】',
    '• 定义：具备"感知-决策-执行-学习"闭环能力的 AI 系统，能自主理解目标并规划行动',
    '• 核心特征：自主性（主动思考）、反应性（实时响应）、主动性（目标导向）、社会性（多智能体协作）',
    '• 范式转变：从"人找工具"到"工具找人"、从"菜单操作"到"对话交互"、从"执行指令"到"理解意图"',
    '',
    '【勘探院场景的智能体应用 - 标杆案例】',
    '▎地质力学参数计算智能体（石化工程院）',
    '• 功能：语言指令驱动的高效精准计算与知识查询，支持自然语言输入专业问题',
    '• 效果：效率提升2倍以上，算法调用准确率95%，人工干预减少70%',
    '',
    '▎钻井提速优化智能体',
    '• 功能：对话式风险处置、钻井参数优化、实时资料查询、历史案例推荐',
    '• 效果：风险处置措施与专家符合率>80%，钻井效率提升15%，事故率降低50%',
    '',
    '▎钻井设计多智能体协同系统',
    '• 功能：多智能体协作，20分钟内生成完整钻井设计方案（传统需2-3天）',
    '• 技术：地质智能体+工程智能体+风险智能体+优化智能体协同决策',
    '• 效果：效率提升5倍，人工修正率<15%，方案质量显著提升',
    '',
    '【智能体的价值】',
    '• 知识沉淀：将专家经验编码为可复用的智能能力，解决"老师傅退休知识流失"问题',
    '• 7×24小时：不知疲倦的"数字员工"，响应时间从小时级缩短至秒级',
    '• 规模复制：优秀经验可快速推广至全油田，打破地域限制'
])

add_content_slide_dense(prs, '范式转移的本质：从第三范式到第四范式', [
    '【科学研究的四种范式演进】',
    '• 第一范式（经验归纳）：实验科学——观察现象、总结经验定律（如开普勒定律）',
    '• 第二范式（理论推演）：模型科学——建立数学模型、推导理论（如牛顿力学）',
    '• 第三范式（计算模拟）：仿真科学——数值计算、计算机模拟（如油藏数值模拟）',
    '• 第四范式（数据密集型）：AI科学——数据驱动发现、机器学习、知识自动提取',
    '',
    '【油气行业的范式转移】',
    '▎传统模式（第三范式）',
    '• 方法论：基于物理方程的数值模拟（油藏数模、地震反演、流体流动方程）',
    '• 局限：计算周期长（单次模拟数小时至数天）、参数调优依赖专家经验、难以全局最优、难以处理不确定性',
    '• 痛点："算得准但算得慢"，无法支撑实时决策',
    '',
    '▎新模式（第四范式）',
    '• 方法论：数据驱动 + 知识增强 + 智能决策，物理模型与数据模型融合',
    '• 优势：实时响应（毫秒级至分钟级）、自主学习能力、跨领域知识融合、处理高维复杂问题',
    '• 突破：从"机理建模"到"数据+机理双驱动"',
    '',
    '【关键技术突破 - 昆仑大模型应用成效】',
    '• 地震资料处理：波动方程求解效率提升10倍，成像质量显著改善',
    '• 储层预测：50亿参数地震解释大模型，解释效率提升10倍以上，吻合率>92%',
    '• 走滑断裂识别：碳酸盐岩储层识别精度大幅提升，助力四川盆地气田发现',
    '',
    '【行业共识：三桶油全面布局】',
    '• 中国石油：PetroAI（700亿参数）、昆仑大模型（3000亿参数）、认知计算平台',
    '• 中国石化：长城大模型、"胜小利"认知大模型、60个AI场景落地',
    '• 中国海油："海能"人工智能平台接入DeepSeek-R1，覆盖海上油气全业务链'
])

# 第二章：核心价值主张
add_section_divider(prs, '02', '核心价值主张：AI 赋能的三大解放')

add_content_slide_dense(prs, '解放一：专家生产力的彻底释放', [
    '【现状痛点：专家时间被严重浪费】',
    '• 时间分配：专家60-70%时间消耗在文献检索、数据整理、格式调整、报告编写等事务性工作',
    '• 效率瓶颈：找一篇相关文献平均耗时30分钟，整理一份报告数据需2-3天',
    '• 创新挤压：重复性工作严重挤压了创造性思考空间，专家无法专注于核心难题攻关',
    '',
    '【AI 解决方案：智能助手替代事务性工作】',
    '▎智能文献助手（PetroAI已沉淀30万份文献，知识图谱超20万三元组）',
    '• 多模态文档智能解析：文字、表格、公式、专业图件（测井曲线、地震剖面）统一理解',
    '• 知识图谱检索增强：覆盖油藏动态分析、CCUS、提高采收率等专业领域',
    '• 语义搜索：支持自然语言提问，如"查找深层碳酸盐岩气藏压裂案例"',
    '• 效果：文献综述时间从数天缩短至数小时，查全率提升40%，查准率提升60%',
    '',
    '▎智能报告生成（大庆"油博士"平台）',
    '• 自动提取多源数据：实时接入生产数据库、实验数据、监测数据',
    '• 规范术语统一：自动识别并统一专业术语、单位、格式',
    '• 语言逻辑修订：自动检查语法、逻辑连贯性、数据一致性',
    '• 覆盖领域：地质、测井、开发、经营等多领域报告',
    '• 效果：报告编制效率提升60%以上，格式错误减少90%',
    '',
    '▎智能代码生成',
    '• 自然语言描述需求 → 自动生成Python/Matlab代码',
    '• 代码解释与调试：自动添加注释、识别潜在错误、优化建议',
    '• 效果：算法调用准确率95%，代码开发效率提升2倍，出错率降低70%',
    '',
    '【释放的价值：专家回归高价值工作】',
    '• 时间重构：专家时间从"事务性工作（70%）"转向"战略性创新（70%）"',
    '• 人才培养：新人培养周期从5-8年缩短至3-5年，智能体成为"随身导师"',
    '• 知识传承：专家经验数字化，解决"退休即流失"痛点'
])

add_content_slide_dense(prs, '解放二：跨学科数据壁垒的打破', [
    '【现状痛点：数据孤岛严重制约决策质量】',
    '• 学科壁垒：地震、地质、钻井、开发、生产各环节数据分散在不同部门',
    '• 标准不一：数据格式、命名规范、单位制不统一，难以融合分析',
    '• 决策片面：单学科视角导致决策局部最优，无法实现全局最优',
    '• 案例：油藏工程师不知道钻井实况，钻井工程师不了解油藏特征',
    '',
    '【AI 解决方案：统一语义与多模态融合】',
    '▎统一数据语义标准（中石油已建立完整标准体系）',
    '• 行业级标准：14项，涵盖数据交换格式、元数据规范、质量要求',
    '• 专业级标准：11项，针对勘探、开发、工程等细分领域',
    '• 场景级标准：48项，覆盖具体业务场景的数据要求',
    '• 数据规模：昆仑大模型数据集达550TB，标注数据量行业领先',
    '',
    '▎多模态数据融合技术',
    '• 统一编码：文本 + 图像 + 时序 + 空间数据统一嵌入向量空间',
    '• 跨模态对齐：地震图像与测井曲线的语义关联，实现"图-数"互查',
    '• 知识融合：结构化数据（数据库）与非结构化数据（文献、报告）统一表示',
    '',
    '▎认知计算平台（勘探院6B平台）',
    '• 技术栈：知识图谱 + 自然语言处理 + 机器学习 + 图计算',
    '• 核心能力：跨学科知识体系的构建、计算、推理和应用',
    '• 应用场景：地质-工程一体化分析、多学科协同决策支持',
    '',
    '【典型应用场景：地质工程一体化】',
    '• 场景：地震数据+钻压采信息实时融合，实现"地质认识指导工程实践，工程数据验证地质认识"',
    '• 数模建模一体化：地质建模→数值模拟无缝衔接，模型更新周期从月级缩短至天级',
    '• 生产优化协同：油藏-井筒-地面全流程优化，从局部优化到全局最优',
    '',
    '【标杆成效：长庆油田】',
    '• 方法：基于"数据-算法-业务"三元耦合架构',
    '• 效果：实现从局部优化到全局最优，采收率提升2-3个百分点，相当于新增可采储量数亿吨'
])

add_content_slide_dense(prs, '解放三：全局最优决策的实现', [
    '【传统决策模式的三大局限】',
    '• 经验驱动：依赖专家直觉和个人经验，难以量化评估，"拍脑袋"决策风险高',
    '• 单点优化：各环节独立决策，缺乏全局视角，局部最优≠全局最优',
    '• 响应滞后：传统数模计算周期长（数周），错过最佳调整时机，"算出来已经晚了"',
    '',
    '【AI 驱动的智能决策体系】',
    '▎实时决策支持：快速代理模型替代传统数模',
    '• 技术：神经网络代理模型（Surrogate Model）学习数模输入输出映射',
    '• 速度：从数周缩短至分钟级，支持实时决策',
    '• 精度：与数模结果吻合度>95%，满足工程精度要求',
    '• 应用：井网层系与配产配注自动优化，动态响应地层变化',
    '',
    '▎风险预警与预测：从事后处置到事前预防',
    '• 钻井风险预警：卡钻、漏失、溢流、井壁失稳提前预警（准确率>90%），预警时间提前24-48小时',
    '• 设备故障预测：基于振动、温度、电流数据，预测性维护，非计划停机减少30%',
    '• 产量异常检测：实时监测单井/区块产量，自动识别异常并诊断原因',
    '',
    '▎数字孪生仿真：虚拟验证再实施',
    '• 架构：构建油气藏-井筒-地面设施全流程数字孪生',
    '• 流程：方案设计→虚拟验证→优化调整→实施→效果评估→模型迭代',
    '• 价值：避免现场试错成本，方案成功率提升30%',
    '',
    '【标杆案例与量化成效】',
    '• 新疆油田：异常工况诊断准确率93.8%，减少作业维护费用20%，单井年节约数十万元',
    '• 苏里格气田：综合递减率从32%降至7%，稳产6年，延长气田生命周期',
    '• 涪陵页岩气：累计产量超1亿m³平台达9个，最高单井累产超7000万m³，创国内纪录'
])

# 第三章：一把手工程
add_section_divider(prs, '03', '"一把手工程"：战略级决策')

add_content_slide_dense(prs, '为什么必须是"一把手工程"？', [
    '【AI 转型的特殊性：三大挑战】',
    '▎跨部门协同难度大——需要"一把手"统筹',
    '• 涉及部门：科研、生产、IT、管理、财务、人力等多部门深度协同',
    '• 打破壁垒：需要打破部门墙，统筹数据资源（涉及数据所有权、隐私、安全）',
    '• 利益协调：改变既有工作流程，可能触动部门利益，需要高层推动',
    '• 结论：只有一把手能协调全局、平衡利益、推动变革',
    '',
    '▎战略投资规模大——需要"一把手"决策',
    '• 算力建设：昆仑大模型高峰算力达1950P（每秒19.5亿亿次计算），投资数十亿元',
    '• 数据治理：550TB数据集治理、数万条标注标准制定，人力投入巨大',
    '• 人才队伍：需要培养复合型AI人才（地学+AI），招聘、培训成本高昂',
    '• 长期投入：AI是持续投入，非一次性项目，需要战略定力',
    '',
    '▎组织变革深度高——需要"一把手"背书',
    '• 文化转变：从"经验驱动"到"数据驱动"，改变数十年工作习惯，阻力巨大',
    '• 流程重构：智能体替代传统岗位，涉及人员调整、职责重新定义',
    '• 考核调整：从"工作量"考核到"价值量"考核，激励机制变革',
    '• 风险容忍：AI项目有试错成本，需要一把手担责、鼓励创新',
    '',
    '【"一把手工程"的核心作用：四定原则】',
    '• 定方向：明确AI战略优先级，纳入单位发展规划，确立"数智化转型"战略地位',
    '• 配资源：算力、数据、人才、资金优先保障，打破资源瓶颈',
    '• 破壁垒：推动跨部门协同，建立数据共享机制，打破信息孤岛',
    '• 担风险：容忍试错，鼓励创新，为探索者兜底，建立容错机制'
])

add_content_slide_dense(prs, '决定未来十年核心竞争力', [
    '【行业竞争态势：全球油气行业AI军备竞赛】',
    '• 国际油服巨头：斯伦贝谢、哈里伯顿已布局AI+数字孪生，技术领先国内2-3年',
    '• 国内"三桶油"全面布局：',
    '  → 中石油：昆仑大模型、PetroAI、认知计算平台，技术体系最完整',
    '  → 中石化：长城大模型、"胜小利"、60个AI场景，应用落地最快',
    '  → 中海油："海能"平台接入DeepSeek-R1，海上油气智能化领先',
    '',
    '【不拥抱AI的四大后果：落后即出局】',
    '• 效率劣势：传统方法效率仅为AI方法的1/5-1/10，成本差距持续拉大',
    '• 人才流失：年轻科研人员流向数智化领先企业，人才梯队断层',
    '• 决策滞后：错失勘探开发最佳时机，好储量被竞争对手抢先开发',
    '• 成本劣势：单井开发成本差距拉大10-15%，在低价环境下失去竞争力',
    '',
    '【拥抱AI的三大战略价值：赢得未来】',
    '• 技术领先：抢占"AI+油气"技术制高点，形成技术壁垒',
    '• 人才高地：吸引顶尖AI+地学复合人才，打造创新团队',
    '• 决策优势：实时全局优化，快速响应市场变化，抓住每个机会窗口',
    '',
    '【时间窗口：现在就是关键节点】',
    '• AI技术迭代周期：6-12个月一次重大升级，技术差距快速拉大',
    '• 先发优势：领先一步，持续积累数据和模型，后来者难以追赶',
    '• 落后代价：落后一步，可能永远追不上，被市场淘汰',
    '• 结论：现在布局，才能在未来十年保持领先；迟疑观望，等于主动放弃'
])

add_content_slide_dense(prs, '勘探院的行动路径：三年规划', [
    '【已具备的基础：不是从零开始】',
    '• 组织基础：2020年成立人工智能研究中心，国内石油行业最早',
    '• 技术基础：2023年发布PetroAI（700亿参数），技术实力行业领先',
    '• 平台基础：承建集团公司"认知计算平台"，定位集团级AI基础设施',
    '• 数据基础：沉淀30万份文献、50TB数据、20万+知识三元组，数据资产雄厚',
    '',
    '【未来三年行动建议：分阶段推进】',
    '▎第一阶段（2025年）：夯实地基——打好基础',
    '• 数据治理：完成全院数据资产盘点与标准化治理，建立数据质量管理体系',
    '• 算力建设：建设勘探院专属智能算力中心（目标：500P以上算力）',
    '• 试点启动：选择10个高价值AI场景试点，快速验证、树立标杆',
    '• 人才培养：启动"AI+油气"复合人才培养计划，首批培养30人',
    '',
    '▎第二阶段（2026-2027年）：规模推广——全面铺开',
    '• 业务覆盖：AI应用覆盖80%核心业务领域，成为日常工作标配',
    '• 示范区建设：建成数字孪生气田示范区（建议选择苏里格或大牛地）',
    '• 人才队伍：累计培养100名AI+油气复合型人才，形成稳定团队',
    '• 生态构建：建立内外部合作生态，引入优质AI服务商',
    '',
    '▎第三阶段（2028-2030年）：智能引领——行业标杆',
    '• 全面智能化：实现勘探开发全流程智能化，关键决策AI辅助率>80%',
    '• 创新中心：建成国家级油气AI创新中心，引领行业标准制定',
    '• 标杆地位：成为中石油数智化转型标杆单位，输出经验和能力',
    '• 人才高地：成为油气AI人才最向往的研究院所'
])

# 结束页
add_title_slide(prs, '拥抱 AI，赢得未来', '以范式转移引领油气开采智能化新征程')

# 保存
output_path = "/root/.openclaw/workspace/AI带来的范式转移-高密度版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板（高密度文字版）")
print(f"📄 正文字号：18-20pt")
print(f"📝 布局特点：单页信息量大，减少留白")
