from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# 创建演示文稿 - 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)
INSPUR_ORANGE = RGBColor(255, 102, 0)
INSPUR_RED = RGBColor(255, 51, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    """添加浪潮装饰元素"""
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    """添加浪潮Logo"""
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    """添加底部标语"""
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    """封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(28)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(20)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_toc_slide(prs, title, items):
    """目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_NAME
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)
    
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_content_slide(prs, title, content_lines):
    """内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = INSPUR_ORANGE
            p.space_before = Pt(10)
            p.space_after = Pt(4)
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆'):
            p.text = '  ' + line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(2)
        elif line == '':
            p.text = ''
            p.space_after = Pt(6)
        else:
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(2)
    
    add_footer_slogan(slide)
    return slide

def add_image_slide(prs, title, image_path, left=Inches(2), top=Inches(1.8), width=Inches(9.333)):
    """图片页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    slide.shapes.add_picture(image_path, left, top, width=width)
    
    add_footer_slogan(slide)
    return slide

def add_table_slide(prs, title, headers, rows):
    """表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5)).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.name = FONT_NAME
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = INSPUR_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.name = FONT_NAME
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容开始 ====================

# P1: 封面
add_title_slide(prs, '油气开采研究所\nAI智能体规划与落地计划', '从"经验驱动"向"数据+知识双轮驱动"转型')

# P2: 目录
add_toc_slide(prs, '目录 CONTENTS', [
    '01  执行摘要与项目背景',
    '02  国际趋势与全球对标',
    '03  国内格局与竞争分析',
    '04  政策环境与标准体系',
    '05  六大核心智能体设计',
    '06  技术架构与数据体系',
    '07  落地路径与实施计划',
    '08  组织保障与资源配置',
    '09  风险管控与竞争策略'
])

# P3: 第一章分隔页
add_section_divider(prs, '01', '执行摘要与项目背景')

# P4: 项目背景
add_content_slide(prs, '项目背景：双重压力与转型需求', [
    '【科研侧痛点】"精度与效率失衡"',
    '• 开发方案设计依赖个人经验，知识传承困难',
    '• 数模建模周期长达数月，专业人才稀缺',
    '• 多专业协同存在壁垒，沟通成本高',
    '',
    '【生产侧痛点】"规模与管控矛盾"',
    '• 低渗透/非常规时代井数激增，管理难度大',
    '• "人盯人"模式难以为继，响应速度慢',
    '• 安全环保风险大，合规压力高',
    '',
    '【现有系统短板】"三多三少"',
    '• 数据多、有效治理少；系统多、跨专业协同少；报表多、智能决策少'
])

# P5: 项目目标与预期收益
add_content_slide(prs, '项目目标与预期收益', [
    '【战略目标】',
    '• 构建"大模型为大脑、智能体为四肢"的协同体系',
    '• 实现科研供给侧与生产需求侧精准对接',
    '• 2027年底完成全面落地',
    '',
    '【核心指标】',
    '• 方案设计周期缩短40%',
    '• 异常诊断准确率>95%',
    '• 核心业务流程智能化覆盖率>80%',
    '• 历史拟合效率提升5倍',
    '• 决策方案论证时间缩短50%',
    '',
    '【战略意义】',
    '• 实现从"经验驱动"向"数据+知识双轮驱动"跨越',
    '• 成为油气行业智能化转型标杆'
])

# P6: 第二章分隔页
add_section_divider(prs, '02', '国际趋势与全球对标')

# P7: 国际巨头AI战略
add_table_slide(prs, '国际巨头AI战略对比', 
    ['国际玩家', '技术路线', '核心场景', '关键数据'],
    [
        ['斯伦贝谢(SLB)', 'NVIDIA NeMo/NIM微服务', '地质勘探、FWI计算', '数字化业务增长9%'],
        ['哈里伯顿', '生成式AI+自治系统', '实时钻井决策、油井测试', '效率提升25%'],
        ['贝克休斯', 'Cordant智能平台', '设备健康管理', '非计划停机减少35%'],
        ['壳牌', 'NVIDIA NeMo底座', '化工Chatbot、3D油藏', '准确率提升30%'],
        ['沙特阿美', 'IBM合作', '井位优化、智能油田', '产量提升35%']
    ]
)

# P8: 国际对标洞察
add_content_slide(prs, '国际对标关键洞察', [
    '【战略模式】底层合作+垂直自研',
    '• 国际大厂普遍采用"底层模型外采/合作 + 垂直Agent自研"策略',
    '• SLB、哈里伯顿等聚焦油气领域专业智能体研发',
    '',
    '【技术路线】云-边-端协同',
    '• 混合云架构：核心地质数据本地存储确保合规',
    '• 计算密集型任务上云，边缘侧实现毫秒级实时响应',
    '',
    '【组织变革】数据科学家+领域专家融合',
    '• 成功企业都建立了跨职能的AI团队',
    '• 数据科学家与地质、钻井工程师深度协作',
    '',
    '【窗口期判断】2-3年技术追赶窗口',
    '• 国际厂商技术领先2-3年',
    '• 但受数据合规限制，无法触碰中国核心生产网',
    '• 需快速建立自有能力壁垒'
])

# P9: 第三章分隔页
add_section_divider(prs, '03', '国内格局与竞争分析')

# P10: 三桶油AI竞赛
add_content_slide(prs, '三桶油"大模型+智能体"竞赛', [
    '【中石油：昆仑生态（行业领先）】',
    '• 技术底座：中国移动+华为盘古+科大讯飞四方共建',
    '• 版本演进：330亿参数→3000亿参数（含800亿多模态）',
    '• 勘探院成效：FWI效率提升10倍，100个场景投产，23类数字员工',
    '',
    '【中石化：时序大模型（差异化）】',
    '• 锚定方向：时序大模型预警预测、下游炼化智能化',
    '• 标杆：镇海炼化"外操无人化、内操智能化"',
    '',
    '【中海油：海洋场景（聚焦）】',
    '• 重点：海上平台无人值守、微电网调度优化',
    '• 差异化：海上风电与油气开采协同降碳'
])

# P11: 供应商布局
add_table_slide(prs, '技术供应商布局',
    ['国内厂商', '核心能力', '石油行业策略', '标杆客户'],
    [
        ['华为', '盘古大模型V5.0、时序增强', '"懂行"+昇腾算力底座', '中石油（昆仑底座）'],
        ['百度', '文心5.0+FM Agent伐谋', '多智能体进化、运筹优化', '国家电网'],
        ['阿里云', '通义千问行业版', '产学研生态、云边协同', '中石化、延长石油'],
        ['科大讯飞', '星火大模型', '语音多模态、数字员工', '中石油（23类数字员工）']
    ]
)

# P12: 第四章分隔页
add_section_divider(prs, '04', '政策环境与标准体系')

# P13: 政策环境
add_content_slide(prs, '政策环境与合规要求', [
    '【国资委"AI+"专项行动（2024）】',
    '• 强制央企建智算中心，催生企业级大模型井喷',
    '• 每个央企至少打造3-5个AI应用标杆场景',
    '',
    '【大模型备案制度】',
    '• 2024年8月，昆仑大模型首批通过备案',
    '• 确立能源合规先发优势',
    '',
    '【数据安全合规】',
    '• 核心数据不出域，确保能源安全',
    '• 建立数据分类分级保护制度',
    '',
    '【行业标准体系】',
    '• 中石油、中石化联合攻坚地震波/时序数据标注标准',
    '• 推动国内标准与国际标准对接'
])

# P14: 第五章分隔页
add_section_divider(prs, '05', '六大核心智能体设计')

# P15: 六大智能体概览
add_table_slide(prs, '六大核心智能体概览',
    ['智能体名称', '服务对象', '核心功能', '预期效果'],
    [
        ['油藏记忆与类比', '科研人员', '知识图谱检索、方案生成', '设计周期缩短40%'],
        ['智能测录井解释', '科研人员', '岩性识别、参数优选', '解释时间减少70%'],
        ['数值模拟加速', '科研人员', '历史拟合辅助、降阶模型', '效率提升5倍'],
        ['生产动态诊断', '生产人员', '异常诊断、治理推荐', '诊断准确率>95%'],
        ['井下作业完整性', '生产人员', '风险预警、压裂优化', '故障漏报率<0.5%'],
        ['开发战术沙盘', '管理层', '策略推演、方案对比', '论证时间缩短50%']
    ]
)

# P16: 插入架构图
add_image_slide(prs, 'AI智能体大脑架构图', '/root/.openclaw/workspace/chart_architecture.png', 
    left=Inches(1), top=Inches(1.5), width=Inches(11.333))

# P17: 科研侧智能体详解
add_content_slide(prs, '科研侧三大智能体详解', [
    '【油藏记忆与类比智能体】',
    '• 构建全球油藏知识图谱，实体类型覆盖油气田、层位、岩性等',
    '• 相似度匹配准确率>90%，自动生成开发方案草案',
    '',
    '【智能测录井解释智能体】',
    '• 多模态数据融合：常规测井、成像测井、核磁测井',
    '• 岩性识别准确率>95%，参数优选与专家吻合度>98%',
    '',
    '【数值模拟加速智能体】',
    '• 强化学习辅助历史拟合，效率提升5倍',
    '• 降阶模型实现秒级预测，误差<5%'
])

# P18: 生产侧智能体详解
add_content_slide(prs, '生产侧三大智能体详解', [
    '【生产动态诊断与优化智能体】',
    '• 实时动态监控单井及油藏动态',
    '• 智能诊断井筒故障，准确率>95%',
    '• 治理措施智能推荐',
    '',
    '【井下作业与完整性智能体】',
    '• 钻完井数据实时分析，井下复杂情况预警',
    '• 井筒完整性监测，故障漏报率<0.5%',
    '• 压裂参数优化，水平井钻遇率提升10%',
    '',
    '【开发战术沙盘智能体】',
    '• 整合生产、财务、市场等多源数据',
    '• 不同开发策略的产量预测与经济效益评估',
    '• 产量配产误差<3%'
])

# P19: 插入对比图
add_image_slide(prs, '传统vs AI模式对比', '/root/.openclaw/workspace/chart_comparison.png',
    left=Inches(1), top=Inches(1.5), width=Inches(11.333))

# P20: 第六章分隔页
add_section_divider(prs, '06', '技术架构与数据体系')

# P21: 技术架构
add_content_slide(prs, '云-边-端三级协同架构', [
    '【L4：交互与决策层】',
    '• 协同研究环境（RDMS）→ 面向科研人员',
    '• 生产调控中心 → 面向生产指挥',
    '• 管理驾驶舱 → 面向管理层',
    '',
    '【L3：智能体层】',
    '• 6大智能体集群 + Agent Orchestration协同机制',
    '',
    '【L2：模型与平台层】',
    '• 昆仑大模型/PetroAI（知识引擎）',
    '• 智能体开发平台（Agent Studio）',
    '',
    '【L1：基础算力与数据层】',
    '• 统一数据湖（地震/钻井/测井/IoT/文献）',
    '• 华为昇腾/鲲鹏算力底座'
])

# P22: 数据飞轮
add_content_slide(prs, '数据飞轮闭环体系', [
    '【数据采集层】',
    '• IoT传感器实时数据采集（压力、温度、流量、振动）',
    '• 业务系统数据同步，外部数据接入',
    '',
    '【数据治理层】',
    '• 数据清洗：缺失值处理、异常值检测',
    '• 数据融合：多源数据对齐、时序对齐',
    '',
    '【数据服务层】',
    '• 特征工程：自动特征提取',
    '• 向量数据库：支持语义检索',
    '• 知识图谱：实体关系抽取',
    '',
    '【数据回流层】',
    '• AI结果数据+人工反馈回流',
    '• 持续微调模型，飞轮效应'
])

# P23: 第七章分隔页
add_section_divider(prs, '07', '落地路径与实施计划')

# P24: 插入路线图
add_image_slide(prs, '三阶段实施路线图', '/root/.openclaw/workspace/chart_roadmap.png',
    left=Inches(1), top=Inches(1.5), width=Inches(11.333))

# P25: 三阶段落地计划
add_content_slide(prs, '三阶段落地计划（2025-2027）', [
    '【第一阶段：夯基垒台（2025）】"让数据说话"',
    '• 2025Q2：数据治理攻坚战，核心数据入湖率90%',
    '• 2025Q3-Q4：油藏记忆、测录井解释智能体1.0上线',
    '',
    '【第二阶段：场景深化（2026）】"让数据思考"',
    '• 2026Q1-Q2：打通科研与生产系统壁垒',
    '• 2026Q2-Q4：生产诊断、井下作业智能体部署，2-3个区块试点',
    '',
    '【第三阶段：融合创新（2027）】"让数据决策"',
    '• 2027Q1-Q2：数值模拟加速与战术沙盘智能体全面应用',
    '• 2027Q4：核心业务流程智能化覆盖率>80%'
])

# P26: 第八章分隔页
add_section_divider(prs, '08', '组织保障与资源配置')

# P27: 插入组织架构图
add_image_slide(prs, '组织保障架构图', '/root/.openclaw/workspace/chart_organization.png',
    left=Inches(1.5), top=Inches(1.5), width=Inches(10.333))

# P28: 预算框架
add_table_slide(prs, '三年预算框架（总计800万）',
    ['预算类别', '2025年', '2026年', '2027年', '合计'],
    [
        ['硬件（算力）', '120万', '80万', '50万', '250万'],
        ['软件（平台）', '80万', '100万', '40万', '220万'],
        ['人力', '90万', '100万', '80万', '270万'],
        ['外部服务', '30万', '20万', '10万', '60万'],
        ['合计', '320万', '300万', '180万', '800万']
    ]
)

# P29: 第九章分隔页
add_section_divider(prs, '09', '风险管控与竞争策略')

# P30: 风险管控
add_table_slide(prs, 'TOP5风险识别与应对',
    ['风险', '概率', '影响', '应对措施'],
    [
        ['数据质量不达标', '高', '高', '建立质量审核机制；数据增强技术'],
        ['技术与业务适配不足', '中', '高', '业务技术联合团队；分阶段测试'],
        ['算力资源不足', '中', '中', '本地+云弹性结合；分布式训练'],
        ['人员接受度低', '中', '中', '分层培训；试点标杆；现场指导'],
        ['项目预算超支', '低', '中', '动态监控；重大投入需审批']
    ]
)

# P31: 差异化优势与竞争策略
add_content_slide(prs, '差异化优势与竞争策略', [
    '【我们的差异化优势】',
    '• 数据资产：30万份文献、50TB数据、20万+知识三元组',
    '• 场景深耕：扎根油气开采研究所，懂科研也懂生产',
    '• 自主可控：基于昆仑大模型底座，核心算法自研',
    '• 快速迭代：miniOPC模式，决策链短，响应速度快',
    '',
    '【竞争策略】',
    '• 短期（1年）：快速完成首批智能体上线，树立标杆',
    '• 中期（1-2年）：深化场景应用，建立技术壁垒',
    '• 长期（2-3年）：构建智能体生态，成为行业标杆',
    '',
    '【合作+竞争关系】',
    '• 与华为：昇腾算力底座合作，掌握核心算法',
    '• 与阿里云：引入补充算力，避免单一供应商绑定'
])

# P32: 下一步行动
add_content_slide(prs, '下一步行动建议', [
    '【即时行动（本周）】',
    '• 成立项目筹备组，明确初步人员分工',
    '• 启动数据资产盘点，识别核心数据源',
    '• 与集团数字化部对接，确认昆仑大模型接入方式',
    '',
    '【短期行动（本月）】',
    '• 完成需求调研，明确6大智能体优先级排序',
    '• 制定详细技术方案，确定与华为/阿里云合作边界',
    '• 编制预算细化方案，申请启动资金',
    '',
    '【中期行动（本季度）】',
    '• 完成数据治理攻坚战首批任务',
    '• 搭建开发测试环境',
    '• 启动油藏记忆与类比智能体MVP开发'
])

# P33: 结束页
add_title_slide(prs, '谢谢', 'AI赋能油气开采  智创能源未来')

# 保存
output_path = "/root/.openclaw/workspace/油气开采研究所AI智能体规划-最终版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
print(f"📷 已插入4张图表")
