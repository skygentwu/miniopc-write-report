#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2025-2026年中国石油勘探开发前沿动态报告PPT
基于浪潮集团模板规范
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 浪潮集团配色方案
INSPUR_BLUE = RGBColor(0, 102, 204)      # 浪潮蓝
INSPUR_ORANGE = RGBColor(255, 102, 0)    # 浪潮橙
TEXT_DARK = RGBColor(51, 51, 51)         # 正文深灰
TEXT_GRAY = RGBColor(102, 102, 102)      # 副标题灰
WHITE = RGBColor(255, 255, 255)

class InspurPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        
    def _add_text(self, slide, left, top, width, height, text, 
                  size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font="微软雅黑"):
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.font.name = font
        return shape
    
    def _add_decorations(self, slide):
        """添加浪潮模板装饰元素"""
        # 右下角橙红色大三角形
        triangle1 = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE, 
            Inches(11.5), Inches(5.8), Inches(2), Inches(1.8)
        )
        triangle1.fill.solid()
        triangle1.fill.fore_color.rgb = INSPUR_ORANGE
        triangle1.line.fill.background()
        triangle1.rotation = 180
        
        # 蓝色小三角形
        triangle2 = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(12.2), Inches(6.3), Inches(1), Inches(0.9)
        )
        triangle2.fill.solid()
        triangle2.fill.fore_color.rgb = INSPUR_BLUE
        triangle2.line.fill.background()
        triangle2.rotation = 180
        
        # 底部标语
        self._add_text(slide, 0.5, 7.0, 12, 0.3, 
                      "未来，因潮澎湃  Inspur in Future",
                      size=10, color=TEXT_GRAY, font="微软雅黑")
    
    def add_cover(self):
        """封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # 左上角Logo占位提示
        self._add_text(slide, 0.5, 0.3, 3, 0.5, "inspur 浪潮", 
                      size=16, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 主标题
        self._add_text(slide, 0.8, 2.5, 11, 1.2,
                      "2025-2026年中国石油勘探开发\n前沿动态报告",
                      size=42, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 副标题
        self._add_text(slide, 0.8, 4.2, 11, 0.8,
                      "深海突破 · 万米深地 · 页岩油革命 · 数智融合",
                      size=24, color=TEXT_GRAY, font="微软雅黑")
        
        # 日期和编制单位
        self._add_text(slide, 0.8, 6.0, 5, 0.4, "2026年3月", 
                      size=14, color=TEXT_GRAY, font="微软雅黑")
        
        self._add_decorations(slide)
        return slide
    
    def add_toc(self):
        """目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        self._add_text(slide, 0.5, 0.3, 3, 0.6, "inspur 浪潮", 
                      size=14, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        self._add_text(slide, 0.8, 1.0, 4, 0.8, "目录", 
                      size=32, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0.8), Inches(1.9), Inches(1.5), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = INSPUR_ORANGE
        line.line.fill.background()
        
        toc_items = [
            ("01", "执行摘要与核心数据"),
            ("02", "深海油气：海洋成为增长极"),
            ("03", "万米深地：向地球深部进军"),
            ("04", "页岩油革命：三大示范区全面建成"),
            ("05", "智能化与数字化转型"),
            ("06", "技术装备迭代升级"),
            ("07", "总结与展望")
        ]
        
        y_pos = 2.3
        for num, title in toc_items:
            self._add_text(slide, 1.0, y_pos, 1, 0.5, num, 
                          size=28, bold=True, color=INSPUR_ORANGE, font="微软雅黑")
            self._add_text(slide, 2.2, y_pos + 0.1, 8, 0.5, title, 
                          size=20, color=TEXT_DARK, font="微软雅黑")
            y_pos += 0.7
        
        self._add_decorations(slide)
        return slide
    
    def add_content_slide(self, title, subtitle=None, bullets=None, table_data=None, highlight=None):
        """内容页通用模板"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # Logo
        self._add_text(slide, 0.5, 0.2, 3, 0.4, "inspur 浪潮", 
                      size=12, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 标题
        self._add_text(slide, 0.5, 0.7, 12, 0.6, title, 
                      size=28, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 副标题
        if subtitle:
            self._add_text(slide, 0.5, 1.3, 12, 0.4, subtitle, 
                          size=16, color=TEXT_GRAY, font="微软雅黑")
        
        # 分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0.5), Inches(1.7), Inches(12), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        line.line.fill.background()
        
        # 要点内容
        if bullets:
            content = "\n".join([f"• {b}" for b in bullets])
            self._add_text(slide, 0.5, 2.0, 12, 4.5, content, 
                          size=16, color=TEXT_DARK, font="微软雅黑")
        
        # 高亮框
        if highlight:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.5), Inches(5.8), Inches(12), Inches(1.0))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(240, 248, 255)
            box.line.color.rgb = INSPUR_BLUE
            box.line.width = Pt(2)
            self._add_text(slide, 0.7, 6.0, 11.5, 0.6, 
                          f"💡 关键洞察：{highlight}", 
                          size=14, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        self._add_decorations(slide)
        return slide
    
    def add_table_slide(self, title, headers, rows, note=None):
        """表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        self._add_text(slide, 0.5, 0.2, 3, 0.4, "inspur 浪潮", 
                      size=12, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        self._add_text(slide, 0.5, 0.7, 12, 0.6, title, 
                      size=28, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        # 创建表格
        num_rows = len(rows) + 1
        num_cols = len(headers)
        
        table = slide.shapes.add_table(num_rows, num_cols, 
                                      Inches(0.5), Inches(1.5), 
                                      Inches(12), Inches(0.6 * num_rows)).table
        
        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = INSPUR_BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = WHITE
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.font.name = "微软雅黑"
        
        # 数据行
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = "微软雅黑"
                # 隔行变色
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 250, 255)
        
        if note:
            self._add_text(slide, 0.5, 6.8, 12, 0.4, note, 
                          size=12, color=TEXT_GRAY, font="微软雅黑")
        
        self._add_decorations(slide)
        return slide
    
    def create_ppt(self):
        """生成完整PPT"""
        # 1. 封面
        self.add_cover()
        
        # 2. 目录
        self.add_toc()
        
        # 3. 执行摘要
        self.add_content_slide(
            "执行摘要：2025年油气勘探里程碑",
            "国家能源局发布全国油气勘探开发十大标志性成果",
            [
                "国内原油产量达2.16亿吨，创历史新高",
                "天然气产量超2600亿立方米，当量首次突破2亿吨",
                "海洋油气产量当量达9000万吨，海上能源堡垒作用凸显",
                "页岩油产量突破850万吨，革命性突破支撑国内增产",
                "万米深地科探井屡破纪录，首次在万米深层发现油气"
            ],
            "2025年是中国油气勘探开发的转折之年，深海、深地、非常规三大方向齐头并进"
        )
        
        # 4. 深海油气 - 成果表格
        self.add_table_slide(
            "深海油气：海洋成为重要增长极",
            ["项目/气田", "2025年产量", "技术特点", "战略意义"],
            [
                ["深海一号", "50亿立方米天然气", "1500米水深、遥控生产", "世界首个远程遥控1500米级半潜平台"],
                ["环海南岛气田群", "1000万吨油当量", "区域协同开发", "较'十三五'末产量翻番"],
                ["渤海油田", "4000万吨油当量", "60余个在产油气田", "全国第一大海上原油基地"]
            ],
            "2025年海洋原油产量突破6600万吨，占全国石油增产量的80%"
        )
        
        # 5. 深海重大发现
        self.add_content_slide(
            "深海勘探：亿吨级重大发现",
            "北部湾盆地潜山领域实现重大突破",
            [
                "惠州19-6油田：我国海上首个深层—超深层碎屑岩亿吨级油田",
                "秦皇岛29-6油田：亿吨级规模，浅层岩性领域重大突破",
                "渤中8-3南气田：渤海海域首个中生界火山岩气田",
                "南海首个深层大中型岩性油田——陆丰13-8油田成功发现",
                "截至2025年三季度末，中国海域获5个新发现，成功评价22个含油气构造"
            ],
            "海洋原油连续5年占全国石油新增产量的60%以上"
        )
        
        # 6. 深海智能化标杆
        self.add_content_slide(
            "深海一号：领航级智能工厂标杆",
            "入选国家首批'领航级智能工厂'培育名单",
            [
                "平台架构：'深海云游'一体化平台 + 数字孪生 + AI算法",
                "覆盖范围：建设-设计-生产-管理全生命周期28个智能化应用模块",
                "年增产天然气6000万立方米，单位产值能耗下降40%",
                "生产效率提升3%，百万吨油气当量用工仅28人（减少50%以上）",
                "世界首个具备远程遥控生产能力的1500米水深半潜式生产储油平台"
            ],
            "深水油气领域数智化转型的典型实践，为行业提供可推广范式"
        )
        
        # 7. 万米深地 - 表格
        self.add_table_slide(
            "万米深地：向地球深部进军",
            ["井名", "完钻井深", "地理位置", "核心成就"],
            [
                ["深地塔科1井", "10,910米", "塔里木盆地", "首次在万米深层发现油气"],
                ["深地川科1井", ">10,000米", "四川盆地", "全球最大尺寸井眼万米深井纪录"],
                ["新安1井", ">10,000米", "准噶尔盆地", "新疆互盈公司万米突破"]
            ],
            "中国工程院院士孙金声：万米深地工程意义比肩探月工程"
        )
        
        # 8. 塔里木盆地主战场
        self.add_content_slide(
            "塔里木盆地：深地工程主战场",
            "超深层已成为我国油气资源增储上产的主阵地",
            [
                "库车山前带：克深20井获高产气流，8个气藏实现千亿方规模增储",
                "富满油田：油气产量当量首次突破460万吨",
                "博孜—大北气田：产量首次突破100亿立方米",
                "顺北超深层油气田：2025年产量当量达324万吨",
                "塔河风化壳下断控穹隆勘探新领域：沙114-2加深井试获百吨高产"
            ],
            "向地球深部挺进是保障我国能源安全的重大战略任务"
        )
        
        # 9. 页岩油革命
        self.add_table_slide(
            "页岩油革命：三大国家级示范区全面建成",
            ["示范区", "核心技术", "年产量", "特点"],
            [
                ["新疆吉木萨尔", "黄金靶体识别、广域支撑压裂等30余项", "170万吨", "产能释放最快"],
                ["大庆古龙", "陆相页岩油高效开发技术", "快速上产中", "老区新资源战略接替"],
                ["胜利济阳", "页岩油精细勘探开发", "持续增产", "技术集成创新示范"]
            ],
            "2025年国内页岩油产量飙升至850万吨以上，支撑全国原油产量达历史最高水平"
        )
        
        # 10. 物探技术突破
        self.add_content_slide(
            "物探技术：2000万道集超大规模采集",
            "中国石油东方物探刷新陆上地震勘探纪录",
            [
                "首次成功应用2000万道集超高密度地震数据采集技术",
                "在准噶尔盆地玛51X井区高质量完成三维地震采集",
                "刷新我国陆上地震勘探采集观测系统强度纪录",
                "为复杂油气藏精准识别提供关键技术支撑",
                "助力页岩油从'资源潜力'向'现实产量'加速转化"
            ],
            "不仅是道数的突破，更是勘探理念与技术模式的跃升"
        )
        
        # 11. 智能化转型
        self.add_content_slide(
            "西南油气田：冲刺500亿立方米年产目标",
            "数字化实践推动效率革命",
            [
                "2025年新井投产达470口，创历史最高纪录，较上年增幅近40%",
                "川中二叠系井平均钻井周期从310天缩短至95.71天",
                "形成175兆帕超高压试油配套技术，创储层改造压力世界纪录",
                "勘探成功率提升至80%，成本降低34%",
                "通过投产井挂图跟踪、井工程全口径滚动等举措确保高效有序"
            ],
            "川渝千亿方天然气生产基地建设提速"
        )
        
        # 12. 技术装备 - 表格1
        self.add_table_slide(
            "技术装备突破（一）：中国石油",
            ["装备/技术", "技术参数", "突破意义"],
            [
                ["EV80可控震源", "8万磅超大吨位、12000米有效信号", "破解超深层勘探信号弱、成像难"],
                ["特深井自动化钻机", "15000米", "国内首套万米钻机"],
                ["万米深层试油技术", "储层改造压力世界纪录", "支撑深地工程实施"]
            ],
            "物探、钻井、压裂、开采全链条核心技术向高端化迈进"
        )
        
        # 13. 技术装备 - 表格2
        self.add_table_slide(
            "技术装备突破（二）：中国石化与中国海油",
            ["企业", "装备/技术", "技术参数", "突破意义"],
            [
                ["中国石化", "应龙科学钻井系统1.0", "钻井全流程智能协同", "国产智能化钻井"],
                ["中国石化", "耐高温混合钻头", "耐245℃", "打破国外技术垄断"],
                ["中国海油", "海洋石油696", "首艘一体化电驱压裂船", "性能参数全球领先"],
                ["中国海油", "3000米水深脐带缆", "CFT光连接器", "工程技术全球前三"]
            ],
            "耐215℃长寿命螺杆钻具、3000米水深连接器实现国产替代"
        )
        
        # 14. 非常规油气
        self.add_content_slide(
            "非常规油气：页岩气勘探新突破",
            "从四川盆地向鄂西地区有效拓展",
            [
                "四川盆地：页岩气产量约270亿立方米，国内首个万亿方深层页岩气储量区",
                "鄂西突破：恩施二叠系页岩层系实现工业气流",
                "新增页岩气地质资源量1329.5亿立方米",
                "中国石化发现国内首个二叠系页岩气大气田",
                "形成的成藏理论认识和高效返排技术具有显著工程实践意义"
            ],
            "非常规天然气产量占比达42%，实现从'常规主导'到'协同发力'转型"
        )
        
        # 15. 核心数据汇总
        self.add_table_slide(
            "2025年核心数据汇总",
            ["指标", "数据", "同比/历史对比"],
            [
                ["国内原油产量", "2.16亿吨", "创历史新高"],
                ["国内天然气产量", ">2600亿立方米", "当量首次突破2亿吨"],
                ["海洋油气产量当量", "9000万吨", "海上能源堡垒凸显"],
                ["页岩油产量", ">850万吨", "革命性突破"],
                ["海洋石油产量", "6800万吨", "占全国增产量80%"],
                ["四川盆地天然气", "800亿立方米", "增量占全国40%以上"]
            ],
            "大力提升油气勘探开发力度'七年行动计划'圆满收官"
        )
        
        # 16. 总结与展望
        self.add_content_slide(
            "总结与展望：四大战略趋势",
            "2026年持续推进'万亿立方米储量、千万吨油田'目标",
            [
                "向深发展：深海+深地双轮驱动，技术极限不断突破",
                "非常规主导：页岩油、页岩气成为增储上产主力",
                "数智融合：AI、数字孪生重塑油气生产模式",
                "装备自主：关键核心技术国产化加速",
                "2026年攻坚重点：页岩油、深层页岩气复杂地质攻关"
            ],
            "油气增储上产从单一赛道冲刺转向多赛道并跑接力"
        )
        
        # 17. 结语
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        self._add_text(slide, 0.5, 0.2, 3, 0.4, "inspur 浪潮", 
                      size=12, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        self._add_text(slide, 0.8, 2.5, 11, 1.2,
                      "从陆地到海洋\n从常规到非常规",
                      size=36, bold=True, color=INSPUR_BLUE, font="微软雅黑")
        
        self._add_text(slide, 0.8, 4.0, 11, 1.0,
                      "油气增储上产不再是单一赛道的冲刺，\n而是多赛道并跑的接力",
                      size=22, color=TEXT_DARK, font="微软雅黑")
        
        self._add_text(slide, 0.8, 5.5, 11, 0.6,
                      "国内油气自主保障能力进一步提升",
                      size=18, bold=True, color=INSPUR_ORANGE, font="微软雅黑")
        
        self._add_decorations(slide)
        
        # 保存
        output_path = "/root/.openclaw/workspace/output/石油勘探前沿动态报告_2025-2026.pptx"
        self.prs.save(output_path)
        print(f"✅ PPT已生成: {output_path}")
        print(f"📊 总页数: {self.slide_count}")
        return output_path

if __name__ == "__main__":
    ppt = InspurPPT()
    ppt.create_ppt()
