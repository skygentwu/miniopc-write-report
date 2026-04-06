from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿 - 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)  # 浪潮蓝
INSPUR_ORANGE = RGBColor(255, 102, 0)  # 浪潮橙
INSPUR_RED = RGBColor(255, 51, 0)  # 装饰三角形用
TEXT_DARK = RGBColor(51, 51, 51)  # 正文深灰
TEXT_GRAY = RGBColor(102, 102, 102)  # 副标题灰

# 字体设置
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    """添加浪潮装饰元素：橙红色大三角形和蓝色小三角形"""
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    """添加浪潮Logo区域（左上角）"""
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    """添加底部标语"""
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    """封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(28)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(20)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_content_slide(prs, title, content_lines):
    """内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = INSPUR_ORANGE
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆'):
            p.text = '  ' + line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
        elif line.startswith('→') or line.startswith('▶'):
            p.text = '    ' + line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_GRAY
            p.space_after = Pt(3)
        elif line == '':
            p.text = ''
            p.space_after = Pt(8)
        else:
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容开始 ====================

# 封面
add_title_slide(prs, 'AI 带来的范式转移', '中石油勘探院油气开采研究所智能化转型战略思考')

# 第一章：从感知到认知
add_section_divider(prs, '01', '从"感知"到"认知"：AI 能力的跃迁')

add_content_slide(prs, '传统 AI vs 现代大模型：能力边界的突破', [
    '【传统 AI：感知智能时代】',
    '• 擅长：图像识别、语音转文字、数据分类',
    '• 局限：单点任务、缺乏推理、无法理解复杂上下文',
    '• 比喻：像一名训练有素的技工，只会特定工序',
    '',
    '【现代大模型：认知智能时代】',
    '▎自主感知：多源数据融合理解',
    '• 同时处理地震数据、测井曲线、生产日志、文献资料',
    '• 案例：昆仑大模型可同时解析文本、图像、时序数据',
    '',
    '▎复杂推理：从"是什么"到"为什么"',
    '• 逻辑链推理：从异常征兆 → 根因分析 → 处置建议',
    '• 案例：钻井风险预警准确率提升至93.8%（新疆油田）',
    '',
    '▎目标规划：从被动响应到主动决策',
    '• 自动生成开发方案、优化井位部署、预测产量趋势',
    '• 案例：大庆油田产量预测准确度达90.74%，效率提升10倍'
])

add_content_slide(prs, '智能体（Agent）：从"工具"到"同事"', [
    '【什么是智能体？】',
    '• 具备"感知-决策-执行-学习"闭环能力的 AI 系统',
    '• 不再是被动等待指令，而是主动理解目标并规划行动',
    '',
    '【勘探院场景的智能体应用】',
    '▎地质力学参数计算智能体（石化工程院）',
    '• 功能：语言指令驱动的高效精准计算与知识查询',
    '• 效果：效率提升2倍以上，算法调用准确率95%',
    '',
    '▎钻井提速优化智能体',
    '• 功能：对话式风险处置、钻井优化、资料查询',
    '• 效果：风险处置措施与专家符合率>80%',
    '',
    '▎钻井设计多智能体',
    '• 功能：20分钟内生成钻井设计整体方案',
    '• 效果：效率提升5倍，人工修正率<15%',
    '',
    '【本质变化】',
    '• 从"人找工具"到"工具找人"',
    '• 从"菜单操作"到"对话交互"'
])

add_content_slide(prs, '范式转移的本质：从第三范式到第四范式', [
    '【科学研究的四种范式】',
    '• 第一范式：经验归纳（实验科学）',
    '• 第二范式：理论推演（模型科学）',
    '• 第三范式：计算模拟（仿真科学）',
    '• 第四范式：数据密集型科学发现（AI科学）',
    '',
    '【油气行业的范式转移】',
    '▎传统模式（第三范式）',
    '• 基于物理方程的数值模拟（油藏数模、地震反演）',
    '• 局限：计算周期长、参数调优依赖经验、难以全局最优',
    '',
    '▎新模式（第四范式）',
    '• 数据驱动 + 知识增强 + 智能决策',
    '• 优势：实时响应、自主学习、跨领域协同',
    '',
    '【关键突破点】',
    '• 地震资料处理：昆仑大模型将波动方程求解效率提升10倍',
    '• 储层预测：50亿参数地震解释大模型，效率提升10倍以上',
    '• 走滑断裂识别：碳酸盐岩储层识别精度大幅提升',
    '',
    '【行业共识】',
    '• 中国石油：PetroAI（700亿参数）、昆仑大模型（3000亿参数）',
    '• 中国石化：长城大模型、"胜小利"认知大模型',
    '• 中国海油："海能"人工智能平台接入DeepSeek-R1'
])

# 第二章：核心价值主张
add_section_divider(prs, '02', '核心价值主张：AI 赋能的三大解放')

add_content_slide(prs, '解放一：专家生产力的彻底释放', [
    '【现状痛点】',
    '• 专家60-70%时间消耗在：文献检索、数据整理、报告编写',
    '• 重复性工作挤压了创造性思考空间',
    '',
    '【AI 解决方案】',
    '▎智能文献助手（PetroAI已沉淀30万份文献）',
    '• 多模态文档智能解析：文字、表格、公式、专业图件',
    '• 知识图谱检索增强：超20万三元组覆盖油藏动态分析、CCUS等',
    '• 效果：文献综述时间从数天缩短至数小时',
    '',
    '▎智能报告生成（大庆"油博士"平台）',
    '• 自动提取多源数据、规范术语、修订语言逻辑',
    '• 覆盖地质、测井、开发、经营多领域',
    '• 效果：报告编制效率提升60%以上',
    '',
    '▎智能代码生成',
    '• 自然语言描述 → 自动生成Python/Matlab代码',
    '• 算法调用准确率95%，效率提升2倍',
    '',
    '【释放的价值】',
    '• 专家时间从"事务性工作"转向"战略性创新"',
    '• 新人培养周期从5-8年缩短至3-5年'
])

add_content_slide(prs, '解放二：跨学科数据壁垒的打破', [
    '【现状痛点】',
    '• 地震、地质、钻井、开发、生产"数据孤岛"',
    '• 学科壁垒导致决策片面，难以全局最优',
    '',
    '【AI 解决方案】',
    '▎统一数据语义标准',
    '• 中石油建立14项行业级、11项专业级、48项场景级数据标准',
    '• 昆仑大模型数据集规模达550TB',
    '',
    '▎多模态数据融合',
    '• 文本 + 图像 + 时序 + 空间数据统一编码',
    '• 跨模态对齐技术实现语义关联',
    '',
    '▎认知计算平台（勘探院6B平台）',
    '• 知识图谱 + 自然语言处理 + 机器学习',
    '• 支撑跨学科知识体系的构建、计算和应用',
    '',
    '【典型应用场景】',
    '• 地质工程一体化：地震数据+钻压采信息实时融合',
    '• 数模建模一体化：地质建模→数值模拟无缝衔接',
    '• 生产优化协同：油藏-井筒-地面全流程优化',
    '',
    '【效果】',
    '• 长庆油田：基于"数据-算法-业务"三元耦合，实现从局部优化到全局最优'
])

add_content_slide(prs, '解放三：全局最优决策的实现', [
    '【传统决策模式】',
    '• 经验驱动：依赖专家直觉，难以量化评估',
    '• 单点优化：各环节独立决策，缺乏全局视角',
    '• 响应滞后：数模计算周期长，错过最佳时机',
    '',
    '【AI 驱动的智能决策】',
    '▎实时决策支持',
    '• 快速代理模型替代传统数模：从数周缩短至分钟级',
    '• 强化学习实现动态优化：井网层系与配产配注自动调整',
    '',
    '▎风险预警与预测',
    '• 钻井风险：卡钻、漏失、溢流提前预警（准确率>90%）',
    '• 设备故障：预测性维护，非计划停机减少30%',
    '',
    '▎数字孪生仿真',
    '• 构建油气藏-井筒-地面全流程数字孪生',
    '• 方案设计→虚拟验证→优化调整→实施',
    '',
    '【标杆案例】',
    '• 新疆油田：异常工况诊断准确率93.8%，减少作业维护费用20%',
    '• 苏里格气田：综合递减率从32%降至7%，稳产6年',
    '• 涪陵页岩气：累计产量超1亿m³平台达9个'
])

# 第三章：一把手工程
add_section_divider(prs, '03', '"一把手工程"：战略级决策')

add_content_slide(prs, '为什么必须是"一把手工程"？', [
    '【AI 转型的特殊性】',
    '▎跨部门协同难度大',
    '• 涉及科研、生产、IT、管理多部门',
    '• 需要打破部门墙，统筹数据资源',
    '• 只有一把手能协调全局',
    '',
    '▎战略投资规模大',
    '• 算力建设：昆仑大模型高峰算力达1950P',
    '• 数据治理：550TB数据集、数万条标注标准',
    '• 人才队伍：需要培养复合型AI人才',
    '',
    '▎组织变革深度高',
    '• 从"经验驱动"到"数据驱动"的文化转变',
    '• 工作流程重构：智能体替代传统岗位',
    '• 考核体系调整：从"工作量"到"价值量"',
    '',
    '【"一把手工程"的核心作用】',
    '• 定方向：明确AI战略优先级',
    '• 配资源：算力、数据、人才、资金',
    '• 破壁垒：推动跨部门协同',
    '• 担风险：容忍试错，鼓励创新'
])

add_content_slide(prs, '决定未来十年核心竞争力', [
    '【行业竞争态势】',
    '• 国际油服巨头：斯伦贝塞、哈里伯顿已布局AI+数字孪生',
    '• 国内"三桶油"：',
    '  → 中石油：昆仑大模型、PetroAI、认知计算平台',
    '  → 中石化：长城大模型、"胜小利"、60个AI场景',
    '  → 中海油："海能"平台接入DeepSeek-R1',
    '',
    '【不拥抱AI的后果】',
    '• 效率劣势：传统方法效率仅为AI方法的1/5-1/10',
    '• 人才流失：年轻科研人员流向数智化领先企业',
    '• 决策滞后：错失勘探开发最佳时机',
    '• 成本劣势：单井开发成本差距拉大10-15%',
    '',
    '【拥抱AI的战略价值】',
    '• 技术领先：抢占"AI+油气"技术制高点',
    '• 人才高地：吸引顶尖AI+地学复合人才',
    '• 决策优势：实时全局优化，快速响应市场',
    '• 成本优势：降本增效，提升竞争力',
    '',
    '【时间窗口】',
    '• AI技术迭代周期：6-12个月',
    '• 落后一步，可能永远追不上',
    '• 现在布局，才能在未来十年保持领先'
])

add_content_slide(prs, '勘探院的行动路径', [
    '【已具备的基础】',
    '• 2020年成立人工智能研究中心',
    '• 2023年发布PetroAI（700亿参数）',
    '• 承建集团公司"认知计算平台"',
    '• 沉淀30万份文献、50TB数据、20万+知识三元组',
    '',
    '【未来三年行动建议】',
    '▎第一阶段（2025）：夯实地基',
    '• 完成全院数据资产盘点与标准化治理',
    '• 建设勘探院专属智能算力中心',
    '• 启动10个高价值AI场景试点',
    '',
    '▎第二阶段（2026-2027）：规模推广',
    '• 覆盖80%核心业务领域',
    '• 建成数字孪生气田示范区',
    '• 培养100名AI+油气复合型人才',
    '',
    '▎第三阶段（2028-2030）：智能引领',
    '• 实现勘探开发全流程智能化',
    '• 建成国家级油气AI创新中心',
    '• 成为中石油数智化转型标杆'
])

# 结束页
add_title_slide(prs, '拥抱 AI，赢得未来', '以范式转移引领油气开采智能化新征程')

# 保存
output_path = "/root/.openclaw/workspace/AI带来的范式转移-浪潮模板版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
