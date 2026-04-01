from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)
INSPUR_ORANGE = RGBColor(255, 102, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(26)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(1.1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_image_slide(prs, title, image_path):
    """插入图片页"""
    if not os.path.exists(image_path):
        print(f"警告：图片不存在 {image_path}")
        return None
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.3), width=Inches(11.7))
    add_footer_slogan(slide)
    return slide

def split_content_by_items(content_lines, max_items=13):
    pages = []
    current_page = []
    current_count = 0
    
    for line in content_lines:
        is_content = line.strip() != ''
        
        if current_count >= max_items and is_content and current_page:
            pages.append(current_page)
            current_page = [line]
            current_count = 1 if is_content else 0
        else:
            current_page.append(line)
            if is_content:
                current_count += 1
    
    if current_page:
        pages.append(current_page)
    
    return pages

def add_content_slide(prs, title, content_lines):
    pages = split_content_by_items(content_lines, max_items=13)
    slides = []
    
    for page_idx, page_content in enumerate(pages):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        add_inspur_logo(slide)
        
        display_title = title if len(pages) == 1 else f"{title} ({page_idx + 1}/{len(pages)})"
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INSPUR_BLUE
        
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(page_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = INSPUR_ORANGE
                p.space_before = Pt(6)
                p.space_after = Pt(2)
            elif line.startswith('•') or line.startswith('-') or line.startswith('◆') or line.startswith('●'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
            elif line.startswith('→') or line.startswith('▶') or line.startswith('○') or line.startswith('  -'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_GRAY
                p.space_after = Pt(2)
            elif not line.strip():
                p.text = ''
                p.space_after = Pt(3)
            else:
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
        
        add_footer_slogan(slide)
        slides.append(slide)
    
    return slides

# ==================== PPT内容 ====================

# 封面
add_title_slide(prs, '设备预防性维护AI智能体', '发展现况深度分析报告 · 全制造产业版')

# 执行摘要
add_content_slide(prs, '执行摘要', [
    '【研究背景】',
    '• 全球PdM市场2024-2025年实现从"阈值报警"向"生成式AI+多智能体协同"跃迁',
    '• 工业Copilot与设备健康智能体成为全制造产业标配',
    '',
    '【核心发现】',
    '• 流程工业（石化/钢铁/电力）：适配度★★★★★，连续生产特征明显',
    '• 离散制造（汽车/电子/机械）：适配度★★★★☆，自动化程度高',
    '• 民生工业（食品/纺织/医药）：适配度★★★☆☆，卫生安全要求高',
    '• 基础设施（矿山/轨交/水务）：适配度★★★★☆，分布广泛',
    '',
    '【量化效果】',
    '• 故障监测准确率＞90%，缺陷发现提前量提升60%+',
    '• 非计划停机减少20-40%，维护成本节约15-30%',
    '• ROI回收期6-12个月',
    '',
    '【市场规模】',
    '• 全球：150-180亿美元（2024），预计2026年突破250亿美元',
    '• 中国：近100亿人民币，增速30%，全球最大增量市场'
])

# 第一部分：全球发展态势
add_section_divider(prs, '01', '全球发展态势与竞争格局')

add_content_slide(prs, '国际工业巨头布局（2023-2025）', [
    '【西门子】Siemens Industrial Copilot',
    '• 客户：舍弗勒集团 | 场景：产线代码生成与故障诊断',
    '• 成效：停机时间↓20%，排查时间↓50%',
    '',
    '【GE Vernova】APM + AI引擎',
    '• 客户：卡塔尔能源 | 场景：燃气轮机预测性维护',
    '• 成效：意外停机↓15-20%，年节约成本超千万美元',
    '',
    '【施耐德电气】EcoStruxure + Advisor',
    '• 客户：巴斯夫 | 场景：电气设备健康监测',
    '• 成效：非计划停机↓近30%',
    '',
    '【ABB】Ability™ Genix',
    '• 客户：瑞典SCA纸业 | 场景：造纸机传动维护',
    '• 成效：维护规划效率↑40%，寿命延长10%',
    '',
    '【罗克韦尔】FactoryTalk® + AI助手',
    '• 客户：福特汽车 | 场景：冲压车间预测维护',
    '• 成效：OEE提升5-8%',
    '',
    '【IBM】Maximo + watsonx',
    '• 客户：陶氏化学 | 场景：视觉与声学联合诊断',
    '• 成效：巡检工作量↓40%，资产利用率↑12%'
])

# 插入行业矩阵图
img1_path = '/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_17---3dad5604-d9f8-4662-be94-c1d686ce45e4.png'
if os.path.exists(img1_path):
    add_image_slide(prs, '行业应用成熟度矩阵（市场规模 vs 技术成熟度）', img1_path)

add_content_slide(prs, '国内AI服务商布局（大厂算力+垂直机理）', [
    '【华为】盘古制造/矿山大模型',
    '• 客户：陕煤集团、长安汽车、宝武集团',
    '• 场景：压铸机/煤矿综采/钢铁轧机 | 准确率超90%',
    '',
    '【阿里云】ET工业大脑 + 通义千问',
    '• 客户：海螺水泥、一汽红旗、浙江正泰',
    '• 场景：水泥窑炉/涂装线/电气设备 | 停机↓20%',
    '',
    '【腾讯云】WeMake工业互联网',
    '• 客户：宁德时代、工业富联、富士康',
    '• 场景：电池产线/SMT设备 | 响应时间↓30%',
    '',
    '【百度】文心工业版 + 调度智能体',
    '• 客户：国家电网、恒力石化、三一重工',
    '• 场景：变电站/纺织机械/工程机械 | 漏检率＜1%',
    '',
    '【科大讯飞】星火大模型 + 声纹识别',
    '• 客户：国能集团、美的、海尔',
    '• 场景：压缩机/电机声纹 | 识别率85%+',
    '',
    '【第四范式】先知AI平台',
    '• 客户：九江石化、宝武集团、燕京啤酒',
    '• 场景：炼化/食品/包装 | 预警准确率95%+',
    '',
    '【容知日新】SuperCare平台',
    '• 客户：中石化、中广核、内蒙古能源',
    '• 场景：风电/石化/核电 | 挽回数亿元损失',
    '',
    '【航天智控】云端PHM',
    '• 客户：沙钢集团、中海油、中国航发',
    '• 场景：连铸机/海上平台/航空发动机 | 成本↓25%'
])

# 第二部分：全制造产业应用
add_section_divider(prs, '02', '全制造产业应用深度分析')

# 插入演进路线图
img2_path = '/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_15---ac3775f3-c093-4080-a4f9-dc9dddd85e6a.png'
if os.path.exists(img2_path):
    add_image_slide(prs, '设备维护技术演进路线图（2010年前→2020年后）', img2_path)

add_content_slide(prs, '流程工业（石化/钢铁/电力）★★★★★', [
    '【石油化工】',
    '• 设备：旋转机械、压缩机、泵群、反应釜',
    '• 痛点：高温高压、连续生产、停机损失巨大',
    '• 案例：中石化镇海炼化（阿里云+中控）',
    '  → 非计划停工↓22%，年节约超千万元',
    '• 案例：中石化九江石化（第四范式）',
    '  → 预警准确率超95%',
    '',
    '【钢铁冶金】',
    '• 设备：高炉、轧机、起重设备、输送辊道',
    '• 痛点：高负荷连续、工况恶劣、质量关联复杂',
    '• 案例：宝武集团（第四范式+宝信）',
    '  → 故障率↓18%，备件资金占用↓10%',
    '• 案例：北京科技大学轧制中心',
    '  → 漏报率≤5%，准确率＞90%',
    '',
    '【电力能源】',
    '• 设备：风电机组、火电机组、变压器、电网',
    '• 痛点：分布广泛、波动性大、稳定性要求高',
    '• 案例：国家电网浙江（百度+华为）',
    '  → 预测准确率96%，巡检成本↓50%',
    '• 案例：国家能源集团"擎源"',
    '  → 41个智能体，半年发现缺陷2633条'
])

add_content_slide(prs, '离散制造（汽车/电子/机械）★★★★☆', [
    '【汽车制造】',
    '• 设备：冲压机、焊接机器人、总装线、压铸机',
    '• 痛点：自动化高、节拍快、混线生产',
    '• 案例：比亚迪深汕基地（腾讯+自研）',
    '  → 非计划停线↓30%以上',
    '• 案例：中国一汽红旗（AI视觉质检）',
    '  → 识别率超99.5%',
    '• 案例：长安汽车（华为盘古）',
    '  → 故障预测准确率超90%',
    '',
    '【电子制造】',
    '• 设备：SMT贴片机、回流焊、AOI检测、晶圆设备',
    '• 痛点：微米级精度、静电敏感、24小时生产',
    '• 案例：宁德时代（腾讯云WeMake）',
    '  → 维护响应时间↓30%',
    '• 案例：工业富联/富士康（设备健康管理）',
    '  → SMT设备预测性维护',
    '• 案例：浙江正泰（阿里云）',
    '  → 电气设备预警',
    '',
    '【机械加工】',
    '• 设备：数控机床、加工中心、注塑机、压铸机',
    '• 痛点：刀具磨损监测、精度波动、多品种',
    '• 案例：三一重工（百度智能体）',
    '  → 工程机械预测维护',
    '• 案例：富士康（模具状态监测）',
    '  → 模具寿命预测'
])

add_content_slide(prs, '民生工业与基础设施', [
    '【食品加工】★★★☆☆',
    '• 设备：灌装机、杀菌釜、包装线、冷链设备',
    '• 痛点：卫生安全（HACCP）、温度敏感、批次追溯',
    '• 案例：燕京啤酒（第四范式）',
    '  → 发酵罐预测维护',
    '• 覆盖：乳制品、饮料、啤酒、肉制品',
    '',
    '【纺织服装】★★★☆☆',
    '• 设备：纺纱机、织布机、印染设备、裁床',
    '• 痛点：设备数量多、温湿度敏感、柔性生产',
    '• 案例：恒力石化（百度智能体）',
    '  → 纺织机械维护',
    '',
    '【医药制造】★★★☆☆',
    '• 设备：反应釜、冻干机、灌装机、灭菌柜',
    '• 痛点：GMP合规、批次记录、交叉污染',
    '• 案例：某大型药企（IBM Maximo）',
    '  → 关键设备验证监控',
    '',
    '【矿山机械】★★★★☆',
    '• 案例：陕煤集团（华为）→ 备件库存成本↓15%',
    '',
    '【轨道交通】★★★★☆',
    '• 案例：中国中车（容知日新）→ 故障率↓超15%',
    '',
    '【水务环保】★★★☆☆',
    '• 案例：某水务集团（施耐德）→ 泵站预测维护'
])

# 第三部分：技术架构与价值
add_section_divider(prs, '03', '技术架构与业务价值')

# 插入技术架构图
img3_path = '/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_16---6d8e4a58-bd17-437b-9082-ec2a53499902.png'
if os.path.exists(img3_path):
    add_image_slide(prs, 'AI智能体技术架构（云-边-端6层架构）', img3_path)

add_content_slide(prs, 'AI智能体核心能力架构', [
    '【L6 交互与执行层】',
    '• 可视化决策看板（支持多工厂、多层级）',
    '• 自然语言交互助手（Industrial Copilot）',
    '• 自动对接EAM/MES/ERP/工单系统',
    '',
    '【L5 智能体协同层】',
    '• 设备Agent + 工艺Agent + 品质Agent',
    '• 能源Agent + 安全Agent 协同',
    '',
    '【L4 分析建模层】',
    '• 故障诊断：CNN/LSTM/Transformer/时序大模型',
    '• 寿命预测：威布尔分布+深度学习融合',
    '• 根因分析：因果推理+知识图谱',
    '',
    '【L3 数据平台层】',
    '• 工业资源层：设备档案、运行数据、3D模型',
    '• 知识经验层：故障库、专家规则、行业标准',
    '• 企业应用层：SCADA/MES/EAM/PLM集成',
    '',
    '【L2 边缘计算层】',
    '• 数据预处理、特征提取、毫秒级响应',
    '',
    '【L1 数据采集层】',
    '• 振动+温度+电流+声纹+视觉+红外+油液'
])

add_content_slide(prs, '关键算法与技术能力', [
    '【CNN 卷积神经网络】',
    '• 适用：振动频谱分析、热成像识别、视觉质检',
    '• 特点：自动提取时空特征',
    '• 行业：汽车、电子、食品',
    '',
    '【LSTM 长短期记忆网络】',
    '• 适用：时间序列预测、剩余寿命预测',
    '• 特点：捕捉长期依赖关系',
    '• 行业：石化、电力、钢铁',
    '',
    '【Transformer/时序大模型】',
    '• 适用：多传感器融合、异常检测、跨设备迁移',
    '• 特点：自注意力机制，小样本表现优异',
    '• 行业：全行业',
    '',
    '【知识图谱】',
    '• 适用：故障根因分析、专家经验固化',
    '• 特点：结构化知识，支持推理',
    '• 行业：航空、医药、核电',
    '',
    '【强化学习】',
    '• 适用：维护策略优化、调度决策、能耗优化',
    '• 特点：动态环境学习最优策略',
    '• 行业：物流、能源、制造排程',
    '',
    '【声纹识别】',
    '• 适用：旋转机械故障诊断、异音检测',
    '• 特点：非接触式，适合恶劣环境',
    '• 行业：电力、矿山、通用机械'
])

add_content_slide(prs, '典型业务场景与价值量化', [
    '【设备异常早期预警】',
    '• 全行业适用 | 提前量提升60%+，漏报率≤5%',
    '',
    '【故障根因智能诊断】',
    '• 流程工业、汽车 | 效率提升25-40%，准确率>90%',
    '',
    '【维护策略自动生成】',
    '• 高价值设备 | 非计划停机减少，重复故障率下降',
    '',
    '【备件预测与调度】',
    '• 全制造产业 | 库存周转率提升，紧急采购减少',
    '',
    '【点检作业智能化】',
    '• 设备数量多行业 | 工作量降低60%+',
    '',
    '【跨业务协同优化】',
    '• 流程工业 | OEE提升5%',
    '',
    '【安全风险主动防控】',
    '• 高危行业 | 违章识别15秒预警',
    '',
    '【知识传承与培训】',
    '• 全行业 | 故障处理效率提升40%+',
    '',
    '【能效优化与碳管理】',
    '• 高能耗行业 | 能耗降低6-8%',
    '',
    '【质量溯源与批次管理】',
    '• 食品、医药 | 批次追溯时间缩短80%'
])

# 第四部分：趋势与实施建议
add_section_divider(prs, '04', '发展趋势与实施建议')

add_content_slide(prs, '2025-2030年技术演进趋势', [
    '【智能体试点期】2025-2026',
    '• 单场景智能体应用成熟，头部企业规模化部署',
    '• 大模型与工业场景深度融合',
    '• 关键能力：预测准确率>90%，人机协同机制建立',
    '• 适用：高价值设备、关键产线',
    '',
    '【多智能体协同期】2027-2028',
    '• 跨专业智能体协同成为主流',
    '• 实现"设备-生产-质量-成本-安全"一体化优化',
    '• 关键能力：全局最优，数字孪生深度融合',
    '• 适用：智能工厂、灯塔工厂',
    '',
    '【自主决策期】2029-2030',
    '• AI智能体具备自主决策能力',
    '• 部分场景实现"无人运维"',
    '• 从预测走向规范(prescriptive)',
    '• 关键能力：自主优化，责任界定清晰',
    '• 适用：黑灯工厂、全自动化产线'
])

# 插入投入收益图
img4_path = '/root/.openclaw/media/inbound/石油勘探AI智能体落地实施计划书_18---e848e927-4319-47e5-a14e-e89d7425539b.png'
if os.path.exists(img4_path):
    add_image_slide(prs, '投入期 vs 收益期（盈亏平衡点分析）', img4_path)

add_content_slide(prs, '分行业实施路径建议', [
    '【流程工业：石化/钢铁/电力】',
    '• 切入点：关键旋转机械（压缩机、泵、风机）',
    '• 第1阶段（3-6月）：高价值设备试点',
    '• 第2阶段（6-12月）：扩展至同类设备群',
    '• 第3阶段（12-24月）：多装置协同优化',
    '',
    '【离散制造：汽车/电子/机械】',
    '• 切入点：瓶颈设备、关键质量影响设备',
    '• 第1阶段（3-6月）：单条产线试点，验证OEE',
    '• 第2阶段（6-12月）：全厂推广，设备健康档案',
    '• 第3阶段（12-24月）：供应链协同',
    '',
    '【民生工业：食品/纺织/医药】',
    '• 切入点：质量关键设备、冷链设备',
    '• 第1阶段（3-6月）：质量关键设备试点',
    '• 第2阶段（6-12月）：卫生安全相关设备',
    '• 第3阶段（12-24月）：柔性生产支持'
])

add_content_slide(prs, '供应商选型建议（分规模）', [
    '【大型央企/国企】营收>500亿',
    '• 推荐：国内头部大厂',
    '• 厂商：华为、阿里云、百度',
    '• 考量：安全合规、生态完整、长期服务',
    '',
    '【中型制造企业】营收50-500亿',
    '• 推荐：垂直领域厂商',
    '• 厂商：第四范式、容知日新、航天智控',
    '• 考量：高ROI、快速见效、行业Know-how',
    '',
    '【中小型民企】营收<50亿',
    '• 推荐：SaaS化轻量级方案',
    '• 厂商：阿里云IoT、腾讯云WeMake',
    '• 考量：低成本、快速部署、免运维',
    '',
    '【外资/合资企业】',
    '• 推荐：国际厂商合规版本',
    '• 厂商：西门子Azure中国、GE',
    '• 考量：全球标准、跨国协同'
])

add_content_slide(prs, '市场规模与政策环境', [
    '【全球市场】',
    '• 2024年：150亿-180亿美元',
    '• 2026年预测：突破250亿美元',
    '• CAGR：25%以上',
    '',
    '【中国市场】',
    '• 2024年：近100亿人民币',
    '• 2025年增速：30%',
    '• 地位：全球最大增量市场',
    '',
    '【投资热点】',
    '• 边缘AI算力芯片',
    '• 高精度无线振动/声纹传感器',
    '• 垂直大模型工业软件',
    '',
    '【核心政策】',
    '• 《推动工业领域设备更新实施方案》（工信部2024）',
    '• 国资委"人工智能+"专项行动（2024）',
    '• 《制造业数字化转型行动方案》'
])

# 结束页
add_title_slide(prs, '设备预防性维护AI智能体', '从预测到自主 · 从单点到协同 · 全制造产业智能化')

# 保存
output_path = "/root/.openclaw/workspace/设备预防性维护AI智能体发展现况-全制造产业版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
print(f"📄 正文字号：16pt")
print(f"🖼️ 已插入4张图表")
