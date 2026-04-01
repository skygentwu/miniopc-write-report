from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)
INSPUR_ORANGE = RGBColor(255, 102, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(26)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(1.1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = INSPUR_ORANGE
            p.space_before = Pt(6)
            p.space_after = Pt(2)
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆'):
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(2)
        elif line.startswith('→') or line.startswith('  -'):
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_GRAY
            p.space_after = Pt(2)
        elif not line.strip():
            p.text = ''
            p.space_after = Pt(3)
        else:
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(2)
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容 ====================

# 封面
add_title_slide(prs, '驱动未来：AI大模型重塑石油勘探', '中石油勘探开发研究院数智化转型战略方案 · 修改版')

# 目录（精简）
add_content_slide(prs, '目录', [
    '【第一部分】人工智能发展新趋势（精简至5大核心洞察）',
    '',
    '【第二部分】石油行业AI应用案例（3个专属案例）',
    '• 案例1：地震资料智能解释与储层预测',
    '• 案例2：智能钻井风险预警与优化',
    '• 案例3：油气藏数值模拟智能加速',
    '',
    '【第三部分】中石油勘探领域专属方案',
    '• 昆仑大模型对接与融合',
    '• 数据资产盘点与治理',
    '',
    '【第四部分】应用场景与实施路径',
    '• 五大AI应用场景（量化ROI）',
    '• 2025-2028分阶段实施路线图'
])

# 第一部分：AI趋势（精简版）
add_section_divider(prs, '01', '人工智能发展新趋势（精简版）')

add_content_slide(prs, '核心洞察1：AI经济影响深远', [
    '【全球预测】',
    '• 2030年AI全球经济影响：22.3万亿美元（占全球GDP 3.7%）',
    '• 数据来源：IDC Macroeconomic Center, March 2025',
    '',
    '【能源行业机遇】',
    '• AI在能源勘探开发领域应用仍处于早期',
    '• 国际油服巨头已全面布局：斯伦贝谢150+AI应用、BP上游可靠性97%',
    '• 中国石油：昆仑大模型启动，行业AI应用拐点已至'
])

add_content_slide(prs, '核心洞察2：从自动化到自主伙伴', [
    '【技术演进】',
    '• 感知AI → 生成式AI → Agentic AI（智能体时代）',
    '• AI正从"自动化工具"升级为"自主业务伙伴"',
    '',
    '【油服企业实践】',
    '• 斯伦贝谢：Lumi数据与AI平台，覆盖钻井、生产、维护全链条',
    '• BP：AI避免10%产能损失，新勘探记录12个发现',
    '• 壳牌：世界经济论坛"数字化转型灯塔工厂"',
    '',
    '【启示】石油行业正从"数字化"迈向"智能化"拐点'
])

add_content_slide(prs, '核心洞察3：大模型+多模态能力跃升', [
    '【模型密集发布】',
    '• 2025-2026年：DeepSeek、豆包、千问、GPT-5等20余款模型迭代',
    '• 趋势：高效稀疏架构、长上下文、多模态融合',
    '',
    '【多模态驱动升级】',
    '• 跨模态语义理解：同时分析文本、图像、时序数据',
    '• 应用：地震资料解释（波形+地质图+文本描述）',
    '• 复杂文档处理：技术报告、标准规范智能理解'
])

add_content_slide(prs, '核心洞察4：智能体(AI Agent)时代', [
    '【定义】智能体 = 大模型+工具调用+自主规划执行',
    '',
    '【演进路径】',
    '• 过去：AI给建议（Chat模式）',
    '• 现在：AI执行步骤、产出文件（Cowork/Manus模式）',
    '',
    '【石油行业应用】',
    '• 地震解释智能体：自动完成从数据加载到成果图件生成',
    '• 钻井优化智能体：实时监控→风险预警→参数调整建议',
    '• 知识管理智能体：文献检索→知识提炼→报告生成'
])

add_content_slide(prs, '核心洞察5：开源生态与成本下降', [
    '【开源平权】',
    '• 开源与闭源模型差距从8%缩小至1.7%（2024-2025）',
    '• 未来格局：开源基础模型+商业服务双轨并行',
    '',
    '【成本快速下降】',
    '• API成本年降9-900倍（Epoch AI估计）',
    '• DeepSeek：0.2元/百万tokens（输入）',
    '• 阿里Qwen：0.5元/百万tokens（输入）',
    '',
    '【启示】企业AI应用门槛大幅降低，规模化落地时机成熟'
])

# 第二部分：石油行业案例（重写）
add_section_divider(prs, '02', '石油行业AI应用案例')

add_content_slide(prs, '案例1：地震资料智能解释与储层预测', [
    '【国际对标：斯伦贝谢Lumi平台】',
    '• 数据处理速度提升：5-10倍',
    '• 地质成像分辨率：更高精度',
    '• 勘探成效：助力BP实现12个新发现',
    '',
    '【技术方案：海岳大模型适配】',
    '• 海岳-70B垂域模型，支持地震数据多模态理解',
    '• 知识图谱：融合地质规律+专家经验',
    '',
    '【建议试点场景】',
    '• 初至波自动拾取、断层智能识别、储层参数预测',
    '→ 预期：解释效率↑60-80%，新井成功率↑15%'
])

add_content_slide(prs, '案例2：智能钻井风险预警与优化', [
    '【国际实践】',
    '• BP：钻井溢流预警检测率90%，大幅减少井下意外',
    '• 斯伦贝谢：150个AI应用覆盖钻井风险预测',
    '',
    '【技术方案】',
    '• 实时数据融合：钻压、转速、扭矩、泥浆参数',
    '• AI预测模型：井壁失稳、井漏、井喷风险早期识别',
    '• 知识图谱：历史事故案例关联分析',
    '',
    '【建议试点场景】',
    '• 复杂井况实时监测（深井、超深井>6000m）',
    '• 钻井参数智能优化（钻速vs风险平衡）',
    '→ 预期：NPT↓25%，事故预警提前30-60分钟'
])

add_content_slide(prs, '案例3：油气藏数值模拟智能加速', [
    '【业务痛点】',
    '• 传统数模：地质建模→数模计算→方案优化，迭代周期数周',
    '• 后果：错过最佳投产时机，"带病"决策增加风险',
    '',
    '【AI解决方案：代理模型】',
    '• AI学习物理方程+历史数据，构建Surrogate Model',
    '• 计算时间：从数周→分钟级',
    '• 实时优化：井网部署、配产配注动态调整',
    '',
    '【对标实践】',
    '• BP：强化学习应用于水力压裂等复杂作业',
    '• 预期：方案迭代效率↑10倍，采收率↑1-3%'
])

# 第三部分：中石油专属方案
add_section_divider(prs, '03', '中石油勘探领域专属方案')

add_content_slide(prs, '昆仑大模型对接与融合', [
    '【昆仑大模型现状】',
    '• 2024年启动建设，能源化工行业首个行业大模型',
    '• 目标：覆盖油气勘探、炼油化工、销售全产业链',
    '',
    '【海岳大模型对接路径】',
    '• 基础设施层：昇腾910B算力适配',
    '• 数据层：勘探院专属知识库与集团数据湖对齐',
    '• 模型层：海岳-70B作为垂域模型补充昆仑通用模型',
    '• 应用层：开发勘探专属智能体',
    '',
    '【具体对接内容】',
    '• 数据标准：SEGY、LAS格式标准化',
    '• 知识融合：海岳石油知识库注入昆仑模型'
])

add_content_slide(prs, '中石油勘探院数据资产盘点', [
    '【数据资产框架】',
    '• 地震数据：XX TB（构造解释、储层预测）',
    '• 测井数据：XX TB（储层评价、产能预测）',
    '• 地质资料：XX 万口井（沉积相分析）',
    '• 生产数据：XX 年历史（开发优化）',
    '• 文献专利：XX 万篇（知识检索）',
    '',
    '【数据治理三步走】',
    'Step 1（3个月）：数据标准化、质量评估、历史数据清洗',
    'Step 2（6个月）：知识图谱构建、专家经验显性化',
    'Step 3（持续）：安全管控、分级分类访问'
])

# 第四部分：应用场景与实施路径
add_section_divider(prs, '04', '应用场景与实施路径')

add_content_slide(prs, '场景1：地震资料智能解释 - ROI测算', [
    '【业务痛点】解释周期长、依赖专家、新人培养慢',
    '',
    '【AI方案】初至波自动拾取+断层识别+储层预测',
    '',
    '【投入产出】',
    '• 投入：180万/年（算力100万+治理50万+实施30万）',
    '• 收益：效率↑70%，人员需求↓50%，新井成功率↑15%',
    '• 投资回收期：12-18个月',
    '',
    '【实施建议】',
    '• 难度：★★★☆☆',
    '• 试点：选成熟油田，对比AI解释vs传统解释效果'
])

add_content_slide(prs, '场景2：钻井风险智能预警 - ROI测算', [
    '【业务痛点】复杂井况事故频发、NPT高、损失大',
    '',
    '【AI方案】实时监测+风险预测+处置建议',
    '',
    '【投入产出】',
    '• 投入：170万（边缘设备50万+模型80万+集成40万）',
    '• 收益：NPT↓25%，事故率↓40%，单井成本↓12%',
    '• 投资回收期：8-12个月',
    '',
    '【实施建议】',
    '• 难度：★★★★☆',
    '• 试点：深井/超深井集中区块先行'
])

add_content_slide(prs, '场景3：数值模拟智能加速 - ROI测算', [
    '【业务痛点】数模计算周期长、方案迭代慢、错过最佳时机',
    '',
    '【AI方案】AI代理模型+实时优化+数字孪生',
    '',
    '【投入产出】',
    '• 投入：240万（算力80万+模型100万+数据60万）',
    '• 收益：计算从天→分钟，迭代效率↑10倍，采收率↑1-3%',
    '• 投资回收期：18-24个月（长期收益显著）',
    '',
    '【实施建议】',
    '• 难度：★★★★★',
    '• 试点：开发中后期油田，验证剩余油预测'
])

add_content_slide(prs, '场景4+5：设备维护+碳排放优化', [
    '【场景4：设备预测性维护】',
    '• 投入：110万 | 收益：寿命↑30%，停机↓35% | 回收期：6-9个月',
    '• 难度：★★☆☆☆ | 建议：电潜泵、抽油机先行',
    '',
    '【场景5：碳排放智能监测（探索性）】',
    '• 投入：120万 | 收益：能耗↓8-12%，绿色油气品牌',
    '• 回收期：12-18个月（含政策收益）',
    '• 难度：★★★☆☆ | 建议：能耗大户区块试点',
    '',
    '【综合评估】',
    '• 短期见效：设备维护、钻井预警（6-12个月回收）',
    '• 中期价值：地震解释、碳排放优化（12-18个月）',
    '• 长期战略：数值模拟加速（18-24个月，收益最大）'
])

add_content_slide(prs, '分阶段实施路线图（2025-2028）', [
    '【2025 试点期】投入300万',
    '• 完成数据资产盘点，启动标准化',
    '• 1-2个场景POC（建议：设备维护+钻井预警）',
    '• 3-6个月验证可行性',
    '',
    '【2026 推广期】投入800万',
    '• 扩展至5-8个场景',
    '• 知识库完善，多场景并行',
    '',
    '【2027 深化期】投入1500万',
    '• 全院推广，AI原生架构',
    '• 跨场景协同，业务重塑',
    '',
    '【2028 成熟期】投入2000万',
    '• 自主智能体协同',
    '• 行业标杆，核心竞争力'
])

add_content_slide(prs, '2025年行动路线图（详细）', [
    '【Q2：基础准备】',
    '• 完成数据资产盘点',
    '• 组建项目团队（院领导挂帅+浪潮+业务专家）',
    '• 确定试点场景',
    '',
    '【Q3：POC开发】',
    '• 启动数据治理',
    '• 试点场景开发',
    '• 3个月交付验证版本',
    '',
    '【Q4：评估推广】',
    '• 评估试点效果',
    '• 制定2026推广计划',
    '• 总结最佳实践'
])

# 结束页
add_title_slide(prs, '谢谢！', '浪潮数字企业有限公司  2026年3月')

# 保存
output_path = "/root/.openclaw/workspace/中石油勘探院AI报告-修改版-40页精华.pptx"
prs.save(output_path)
print(f"✅ 修改版PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 浪潮模板")
print(f"\n修改要点：")
print("1. 精简趋势：10个洞察→5个核心观点（删除冗余）")
print("2. 重写案例：3个石油行业案例（替代建筑/化工/矿山）")
print("3. 中石油专属：昆仑大模型对接+数据资产盘点")
print("4. 量化ROI：5个场景投入产出+回收期")
print("5. 实施路线：2025-2028分阶段规划")
