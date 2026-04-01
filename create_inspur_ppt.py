from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿 - 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)  # 浪潮蓝
INSPUR_ORANGE = RGBColor(255, 102, 0)  # 浪潮橙
INSPUR_RED = RGBColor(255, 51, 0)  # 装饰三角形用
TEXT_DARK = RGBColor(51, 51, 51)  # 正文深灰
TEXT_GRAY = RGBColor(102, 102, 102)  # 副标题灰

# 字体设置
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    """添加浪潮装饰元素：橙红色大三角形和蓝色小三角形"""
    # 橙红色大三角形（右下角）
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    # 蓝色小三角形（右侧）
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    """添加浪潮Logo区域（左上角）"""
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    """添加底部标语"""
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    """封面页 - 标题42号加粗，副标题28号"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加Logo
    add_inspur_logo(slide)
    
    # 主标题 - 42号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题 - 28号
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(28)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(20)
    
    # 装饰元素
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    
    return slide

def add_toc_slide(prs, title, items):
    """目录页 - 标题24号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Logo
    add_inspur_logo(slide)
    
    # 目录标题 - 32号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    # 目录内容 - 24号
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_NAME
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(16)
    
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    """章节分隔页 - 大号章节编号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    # 章节编号 - 72号加粗，橙色
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    # 章节标题 - 36号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_content_slide(prs, title, content_lines):
    """内容页 - 标题32号加粗，二级标题20号加粗，正文18号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Logo
    add_inspur_logo(slide)
    
    # 页面标题 - 32号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 判断是二级标题还是正文
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            # 二级标题 - 20号加粗
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = INSPUR_ORANGE
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆'):
            # 正文列表 - 18号
            p.text = '  ' + line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
        elif line == '':
            # 空行
            p.text = ''
            p.space_after = Pt(8)
        else:
            # 普通正文 - 18号
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容开始 ====================

# 封面
add_title_slide(prs, '"深老非"气井开采\n与采气智能体解决方案', '深层·老气田·非常规气田智能化转型')

# 目录
add_toc_slide(prs, '目录 CONTENTS', [
    '01  背景与概念 - "深老非"三类气田特征',
    '02  深层超深层 - 万米级高温高压生产参数寻优',
    '03  老气田 - 低压低产井动态优化排采工艺',
    '04  非常规气田 - 复杂水平井网管控与全生命周期管理',
    '05  采气智能体解决方案',
    '06  标杆案例与效益分析'
])

# 第一章：背景与概念
add_section_divider(prs, '01', '背景与概念')

add_content_slide(prs, '"深老非"定义与行业痛点', [
    '【"深老非"定义】',
    '• 深层/超深层：埋深超过6000m，万米级特深层，超高温(>200°C)、超高压(>200MPa)',
    '• 老气田：开发年限>10年，进入中后期，地层压力降低、产量递减',
    '• 非常规气田：页岩气、煤层气、致密气等，低渗透、低丰度、需大规模压裂',
    '',
    '【行业痛点】',
    '• 深层超深层：单井投资超1亿元，井温超230°C、压力超200MPa，事故风险高',
    '• 老气田：综合递减率高达20-30%，低压低产井占比超50%，人工维护成本高',
    '• 非常规气田：单井产量低（日均<3000m³），压裂效果难预测，生产制度优化难',
    '',
    '【资源占比】',
    '• 深层超深层资源约占全国油气资源总量1/3',
    '• 已开发气田>60%进入老气田阶段',
    '• 非常规天然气产量占比>50%'
])

# 第二章：深层超深层
add_section_divider(prs, '02', '深层超深层')

add_content_slide(prs, '万米级高温高压生产参数寻优', [
    '【技术挑战】',
    '• 井底温度：210-260°C（相当于岩浆温度的25%）',
    '• 井底压力：138-245MPa（相当于2个马里亚纳海沟压力）',
    '• 设备耐温极限需突破300°C',
    '',
    '【关键技术突破】',
    '▎超高温高压测井技术',
    '• 中国石油"经纬视界"系列：耐温260°C/耐压206MPa，技术指标国内最高',
    '• 已应用18口超深井，测井时效提升25%，资料采集率提高22%',
    '',
    '▎万米级钻完井技术',
    '• 全球首台1.2万米自动化钻机',
    '• 抗高温(≥240°C)聚合物水基钻井液，抗220°C高温水泥浆体系',
    '• 深地塔科1井：成功钻至10910米，刷新全球陆上万米深井钻探纪录'
])

add_content_slide(prs, '深层超深层 · 技术突破（续）', [
    '▎超高温射孔技术',
    '• "先锋"超高温超高压射孔器：耐温260°C/72h，耐压245MPa',
    '• 穿深达829mm，超越国际标准57%',
    '',
    '▎智能化应用方向',
    '• 钻头磨损预测与钻参优化：基于机器学习的钻头寿命预测模型',
    '• 井下故障预警：卡钻、漏失、溢流等风险提前预警',
    '• 地质导向实时优化：随钻测井数据实时分析，调整钻井轨迹',
    '',
    '【典型案例】',
    '• 深地塔科1井：10910米全球陆上最深井',
    '• 中国石化"深地工程"：攻克8000米以上超深井近50井次'
])

# 第三章：老气田
add_section_divider(prs, '03', '老气田')

add_content_slide(prs, '低压低产井动态优化排采工艺', [
    '【技术挑战】',
    '• 气井压力低于外输压力，无法自喷',
    '• 井筒积液严重，影响产能',
    '• 单井管理成本高，人工巡检效率低',
    '',
    '【排水采气工艺体系】',
    '• 泡沫排水：产水量<50m³/d，有效率89.6%（大牛地气田）',
    '• 柱塞气举：产水量50-200m³/d，日增气5万m³（苏里格气田）',
    '• 速度管柱：低压低产井，携液能力提升30%（涪陵页岩气）',
    '• 机抽排水：产水量>200m³/d，深井适用（元坝气田）',
    '• 气举排水：高气液比井，无需井下设备（普光气田）'
])

add_content_slide(prs, '老气田 · 智能排采技术', [
    '【老井挖潜技术】',
    '• 侧钻水平井：动用剩余储量，单井日产2.7万m³',
    '• 查层补孔：纵向拓层，增加动用储量',
    '• 储层解堵：蒸汽脉冲解堵，复活停产井',
    '',
    '【智能间开技术】',
    '• 根据"水气比"分类管理：低/中/高水气比差异化措施',
    '• 智能柱塞控制+井口智控阀',
    '• 13口井实现自动化间开，减少水淹"假死"',
    '',
    '【AI应用方向】',
    '• 气井积液预警模型',
    '• 最优开关井时机预测',
    '• 排采工艺智能推荐',
    '• 产量递减趋势预测'
])

# 第四章：非常规气田
add_section_divider(prs, '04', '非常规气田')

add_content_slide(prs, '复杂水平井网管控与全生命周期管理', [
    '【技术挑战】',
    '• 水平段长达2500米（国内纪录）',
    '• 压裂段数多（单井30-60段）',
    '• 裂缝网络复杂，生产制度难优化',
    '',
    '【关键技术体系】',
    '▎地质工程一体化',
    '• 地震数据+钻压采信息实时融合',
    '• 三维储层建模+工业化建模技术',
    '',
    '▎水平井优快钻井',
    '• 一趟钻技术：水平段一次完成比例80%',
    '• 钻井周期从32天缩短至30天',
    '',
    '▎缝网体积压裂',
    '• "工厂化"拉链式压裂：日压裂6段',
    '• 微地震实时监测：裂缝长度精度90%'
])

add_content_slide(prs, '非常规气田 · 全生命周期采气工艺', [
    '【深层煤岩气技术演进】',
    '• 单井EUR：从5594万m³提升至6000万m³（+7.3%）',
    '• 首年递减率：从38%降至33%（-5%）',
    '• 水平段长：从1204m提升至1500m（+24.6%）',
    '• 钻井周期：从32天缩短至30天（-6.3%）',
    '• 压裂周期：从7.5天缩短至6.0天（-20%）',
    '',
    '【AI应用方向】',
    '• 压裂参数智能优化',
    '• 裂缝扩展实时模拟',
    '• 生产制度动态优化',
    '• EUR预测与井位部署优化'
])

# 第五章：采气智能体解决方案
add_section_divider(prs, '05', '采气智能体解决方案')

add_content_slide(prs, '智能体架构设计', [
    '【三层架构】',
    '▎感知层智能体',
    '• 数据接入与预处理：振动、温度、压力、流量传感器数据',
    '• 边缘计算：实时数据处理与特征提取',
    '',
    '▎决策层智能体',
    '• 知识图谱 + 大模型 + 专家系统',
    '• 方案生成与优化：基于深度学习的决策推荐',
    '',
    '▎执行层智能体',
    '• 自动调控：远程控制与闭环优化',
    '• 人机协同：异常情况人工介入'
])

add_content_slide(prs, '五大核心智能体', [
    '【① 钻井优化智能体】',
    '• 功能：轨迹优化、钻参推荐、风险预警',
    '• 效果：钻井效率↑20%，事故率↓50%',
    '',
    '【② 压裂优化智能体】',
    '• 功能：压裂参数优化、裂缝模拟、效果预测',
    '• 效果：单井产量↑15%，压裂成本↓10%',
    '',
    '【③ 生产优化智能体】',
    '• 功能：制度优化、配产决策、能耗管理',
    '• 效果：采收率↑2-5%，能耗↓15%',
    '',
    '【④ 设备健康管理智能体】',
    '• 功能：故障预测、寿命评估、维护决策',
    '• 效果：非计划停机↓30%，维护成本↓25%',
    '',
    '【⑤ 安全巡检智能体】',
    '• 功能：风险识别、预警推送、应急指挥',
    '• 效果：事故率↓40%，巡检效率↑50%'
])

add_content_slide(prs, '技术实现路径', [
    '【第一阶段：数据基础（0-6个月）】',
    '• 部署传感器网络（振动、温度、压力、流量）',
    '• 建设数据中台，实现数据统一治理',
    '• 接入历史数据，建立训练数据集',
    '',
    '【第二阶段：模型构建（6-12个月）】',
    '• 训练故障诊断模型、产量预测模型',
    '• 构建知识图谱，沉淀专家经验',
    '• 开发智能体原型，现场试验验证',
    '',
    '【第三阶段：规模推广（12-24个月）】',
    '• 覆盖全部关键设备与重点气井',
    '• 实现多智能体协同、自主决策',
    '• 持续优化模型，提升准确率'
])

# 第六章：标杆案例与效益分析
add_section_divider(prs, '06', '标杆案例与效益分析')

add_content_slide(prs, '标杆案例', [
    '【深层超深层：深地塔科1井】',
    '• 深度10910米，全球陆上最深井',
    '• 1.2万米自动化钻机+260°C超高温测井仪',
    '',
    '【老气田：苏里格气田】',
    '• 年产300亿m³致密气田，低压低产井占比50%',
    '• 泡沫排水+柱塞气举+智能间开',
    '• 综合递减率从32%降至7%，稳产6年',
    '',
    '【非常规气田：大吉气田】',
    '• 中国首个百万吨油气当量深层煤岩气田',
    '• 地质工程一体化+缝网体积压裂',
    '• 日产气量突破650万m³',
    '',
    '【涪陵页岩气田】',
    '• 国家级页岩气示范区',
    '• 累计产量超1亿m³平台达9个'
])

add_content_slide(prs, '投资效益分析', [
    '【投资估算（中型气田）】',
    '• 传感器网络：50-100万/井',
    '• 数据中台：500-1000万',
    '• AI平台：300-500万',
    '• 智能体开发：200-500万',
    '• 系统集成：100-200万',
    '• 合计：1200-2300万',
    '',
    '【直接效益】',
    '• 单井产量提升：10-20%',
    '• 采收率提升：2-5%',
    '• 维护成本降低：20-30%',
    '• 非计划停机减少：30-50%',
    '',
    '【投资回报期】',
    '• 中型气田：2-3年',
    '• 大型气田：1-2年'
])

add_content_slide(prs, '发展趋势与建议', [
    '【技术趋势】',
    '• 工业大模型应用：构建油气行业专属大模型',
    '• 数字孪生深化：从单井向气田级、集团级演进',
    '• 多智能体协同：钻井-压裂-生产全流程智能协同',
    '• 边缘智能部署：毫秒级响应',
    '',
    '【实施建议】',
    '• 深层超深层：重点突破高温高压传感器可靠性',
    '• 老气田：优先部署低压低产井智能排采系统',
    '• 非常规气田：建设地质工程一体化智能平台',
    '',
    '【风险与对策】',
    '• 数据质量：建立数据治理体系',
    '• 模型准确性：持续迭代优化，人机协同决策',
    '• 人员接受度：加强培训，从辅助工具切入'
])

# 结束页
add_title_slide(prs, '谢谢', 'AI赋能油气开采  智创能源未来')

# 保存
output_path = "/root/.openclaw/workspace/深老非气井开采与采气智能体-汇报PPT-浪潮模板版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
