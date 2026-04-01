from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)
INSPUR_ORANGE = RGBColor(255, 102, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(26)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_section_divider(prs, section_num, section_title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(2), Inches(1.1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def split_content_by_items(content_lines, max_items=13):
    """按条目数分页，每页最多13个有效条目"""
    pages = []
    current_page = []
    current_count = 0
    
    for line in content_lines:
        is_content = line.strip() != ''
        
        if current_count >= max_items and is_content and current_page:
            pages.append(current_page)
            current_page = [line]
            current_count = 1 if is_content else 0
        else:
            current_page.append(line)
            if is_content:
                current_count += 1
    
    if current_page:
        pages.append(current_page)
    
    return pages

def add_content_slide(prs, title, content_lines):
    """内容页 - 高密度版本"""
    pages = split_content_by_items(content_lines, max_items=13)
    slides = []
    
    for page_idx, page_content in enumerate(pages):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        add_inspur_logo(slide)
        
        display_title = title if len(pages) == 1 else f"{title} ({page_idx + 1}/{len(pages)})"
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INSPUR_BLUE
        
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(page_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = INSPUR_ORANGE
                p.space_before = Pt(6)
                p.space_after = Pt(2)
            elif line.startswith('•') or line.startswith('-') or line.startswith('◆') or line.startswith('●'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
            elif line.startswith('→') or line.startswith('▶') or line.startswith('○') or line.startswith('  -'):
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_GRAY
                p.space_after = Pt(2)
            elif not line.strip():
                p.text = ''
                p.space_after = Pt(3)
            else:
                p.text = line
                p.font.name = FONT_NAME
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(2)
        
        add_footer_slogan(slide)
        slides.append(slide)
    
    return slides

def add_image_slide(prs, title, image_path):
    """插入图片页 - 用于展示豆包生成的图表"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_inspur_logo(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    # 插入图片 - 居中，占满大部分空间
    slide.shapes.add_picture(
        image_path,
        Inches(0.8), Inches(1.3),
        width=Inches(11.7)
    )
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容 ====================

# 封面
add_title_slide(prs, '油气开采研究所', 'AI智能体规划与落地计划')

# 执行摘要
add_content_slide(prs, '执行摘要', [
    '【项目背景】',
    '• 科研侧痛点："精度与效率失衡"——方案设计依赖经验，数模周期长达数月',
    '• 生产侧痛点："规模与管控矛盾"——低渗透/非常规时代井数激增，传统模式难以为继',
    '• 系统短板："三多三少"——数据多治理少、系统多协同少、报表多决策少',
    '',
    '【核心目标】',
    '• 构建"大模型为大脑、智能体为四肢"的协同体系',
    '• 实现科研供给侧与生产需求侧精准对接',
    '• 2027年底完成全面落地，核心业务流程智能化覆盖率超80%',
    '',
    '【预期收益】',
    '• 方案设计周期缩短40%，异常诊断准确率>95%',
    '• 实现从"经验驱动"向"数据+知识双轮驱动"跨越'
])

# 插入豆包生成的对比图
add_image_slide(prs, '传统模式 vs AI智能体模式对比', 'chart_comparison.png')

# 第一章：行业趋势
add_section_divider(prs, '01', '行业趋势与对标分析')

add_content_slide(prs, '国际趋势：从预测性AI到生成式AI Agent', [
    '【国际油服巨头】',
    '• 斯伦贝谢(SLB)：联合NVIDIA，基于NeMo/NIM微服务，2024年数字化业务增长9%',
    '• 哈里伯顿：生成式AI+自治系统，已在非传统油气井商业化测试',
    '• 贝克休斯：AI驱动设备健康大模型，深度集成企业级软件',
    '',
    '【国际石油大厂】',
    '• 主流策略："底层模型外采/合作 + 垂直Agent自研"',
    '• 壳牌：NVIDIA NeMo底座，准确率提升30%，训练时间缩短20%',
    '• 沙特阿美：IBM合作，井位优化，产量提升35%',
    '',
    '【关键洞察】',
    '• 平台化合作特征明显，巨头依托底层算力构建领域智能体',
    '• 技术路线趋同：大模型底座+垂直场景Agent'
])

add_content_slide(prs, '国内格局：三桶油的"大模型+智能体"竞赛', [
    '【中石油：昆仑生态（领先）】',
    '• 技术底座：中国移动+华为盘古+科大讯飞四方共建',
    '• 版本演进：2024年8月330亿参数→2025年5月3000亿参数',
    '• 勘探院成效：全波形反演效率提升10倍，100个场景投产，23类数字员工',
    '',
    '【中石化：时序大模型（差异化）】',
    '• 锚定方向：时序大模型预警预测、下游炼化智能化',
    '• 标杆：镇海炼化"外操无人化、内操智能化"',
    '',
    '【中海油：海洋场景（聚焦）】',
    '• 重点：海上平台无人值守、微电网调度优化',
    '• 差异化：海上风电与油气开采协同降碳'
])

add_content_slide(prs, '技术供应商布局与政策环境', [
    '【国内AI服务商】',
    '• 华为：盘古V5.0+昇腾算力，"懂行"+时序增强，中石油昆仑底座',
    '• 百度：文心5.0+FM Agent伐谋，多智能体进化，国家电网光明大模型',
    '• 阿里云：通义千问行业版，产学研生态，云边协同',
    '• 科大讯飞：星火大模型，语音多模态+数字员工',
    '',
    '【国际厂商】',
    '• Palantir、C3.ai基本退出中国核心生产网',
    '• 微软仅通过Azure提供文档/研发辅助，无法触碰核心生产数据',
    '',
    '【政策环境】',
    '• 国资委"AI+"专项行动（2024）：强制央企建智算中心',
    '• 昆仑大模型备案：2024年8月首批通过，确立能源合规先发优势',
    '• 行业标准：中石油、中石化联合攻坚地震波/时序数据标注标准'
])

# 第二章：智能体体系
add_section_divider(prs, '02', '智能体体系设计')

add_content_slide(prs, '需求洞察：科研侧+生产侧双重压力', [
    '【科研侧痛点】"精度与效率失衡"',
    '• 开发方案设计高度依赖个人经验',
    '• 数模建模周期长达数月',
    '• 多专业协同存在壁垒',
    '',
    '【生产侧痛点】"规模与管控矛盾"',
    '• 低渗透/非常规时代井数激增',
    '• "人盯人"模式难以为继',
    '• 安全环保风险大',
    '',
    '【现有系统短板】"三多三少"',
    '• 数据多、有效治理少',
    '• 系统多、跨专业协同少',
    '• 报表多、智能决策少',
    '• 核心问题：系统缺乏"认知能力"'
])

add_content_slide(prs, '六大核心智能体（一）', [
    '【油藏记忆与类比智能体】',
    '• 服务对象：科研人员',
    '• 核心功能：基于知识图谱检索全球类似油藏，自动生成开发方案草案',
    '• 预期效果：方案设计周期缩短40%，相似油藏推荐准确率>90%',
    '',
    '【智能测录井解释智能体】',
    '• 服务对象：科研人员',
    '• 核心功能：多模态大模型自动识别岩性、流体，优选处理解释参数',
    '• 预期效果：解释时间减少70%，参数优选与专家吻合度>98%',
    '',
    '【数值模拟加速智能体】',
    '• 服务对象：科研人员',
    '• 核心功能：辅助历史拟合，自动调整参数；生成降阶模型秒级预测',
    '• 预期效果：历史拟合效率提升5倍，预测时间从小时级降至分钟级'
])

add_content_slide(prs, '六大核心智能体（二）', [
    '【生产动态诊断与优化智能体】',
    '• 服务对象：生产人员',
    '• 核心功能：实时监控单井及油藏动态，智能诊断异常并推荐治理措施',
    '• 预期效果：异常诊断准确率>95%，油藏预测准确率98.7%',
    '',
    '【井下作业与完整性智能体】',
    '• 服务对象：生产人员',
    '• 核心功能：分析钻完井及生产数据，预警井筒完整性风险，优化压裂参数',
    '• 预期效果：井下故障漏报率<0.5%，水平井钻遇率提升10%',
    '',
    '【开发战术沙盘智能体】',
    '• 服务对象：管理层',
    '• 核心功能：基于实时数据与大模型推演，对比不同开发策略的产量与效益',
    '• 预期效果：决策方案论证时间缩短50%，产量配产误差<3%'
])

# 第三章：技术架构
add_section_divider(prs, '03', '技术架构')

add_content_slide(prs, '整体架构：云-边-端三级协同', [
    '【L4：交互与决策层（应用）】',
    '• 协同研究环境（RDMS）→ 面向科研人员',
    '• 生产调控中心 → 面向生产指挥',
    '',
    '【L3：智能体层（执行者）】',
    '• 6大智能体集群 + Agent Orchestration协同机制',
    '',
    '【L2：模型与平台层（大脑）】',
    '• 昆仑大模型/PetroAI（知识引擎）',
    '• 智能体开发平台（Agent Studio）',
    '',
    '【L1：基础算力与数据层（矿藏）】',
    '• 统一数据湖（地震/钻井/测井/IoT/文献）',
    '• 数据治理工具 + 华为昇腾/鲲鹏算力底座'
])

add_content_slide(prs, '数据流与大模型集成', [
    '【数据飞轮闭环】',
    '• 数据源 → 数据湖（IoT实时数据+静态地质数据汇聚）',
    '• 数据湖 → 智能体（RAG技术向量化供大模型调用）',
    '• 智能体 → 决策输出（调用数模软件→生成优化参数）',
    '• 决策输出 → 数据湖（AI结果+人工反馈回流→持续微调）',
    '',
    '【与大模型集成："1+N"模式】',
    '• 能力调用：智能体作为"插件"注册到大模型平台，大模型作"路由器"',
    '• 能力增强：利用大模型通用能力（代码生成、逻辑推理）增强智能体',
    '• 私有化部署：核心地质数据不出域，本地昇腾推理确保安全'
])

# 插入豆包生成的架构图
add_image_slide(prs, '技术架构图（AI智能体大脑+分层体系）', 'chart_architecture.png')

# 第四章：落地路径
add_section_divider(prs, '04', '落地路径（2025-2027）')

add_content_slide(prs, '第一阶段：夯基垒台（2025年）', [
    '【主题】"让数据说话"',
    '',
    '【关键里程碑】',
    '• 2025Q1-Q2：数据治理攻坚战，核心数据入湖率90%',
    '• 2025Q2：搭建智能体开发基础设施',
    '• 2025Q3：上线油藏记忆与类比智能体1.0',
    '• 2025Q4：上线智能测录井解释智能体1.0',
    '',
    '【风险应对】',
    '• 风险：数据质量差、业务专家不配合',
    '• 应对：一把手工程纳入绩效考核；业务+IT混编团队'
])

add_content_slide(prs, '第二阶段：场景深化（2026年）', [
    '【主题】"让数据思考"',
    '',
    '【关键里程碑】',
    '• 2026Q1-Q2：打通科研与生产系统壁垒，动静态数据实时联动',
    '• 2026Q2-Q3：部署生产动态诊断与井下作业智能体',
    '• 2026Q3-Q4：在2-3个重点区块（页岩油/致密气）试点',
    '• 2026Q4：重大安全隐患智能识别准确率>94%',
    '',
    '【风险应对】',
    '• 风险：模型误报导致"狼来了"心理',
    '• 应对：人机协同复核机制；AI建议+人工确认'
])

add_content_slide(prs, '第三阶段：融合创新（2027年）', [
    '【主题】"让数据决策"',
    '',
    '【关键里程碑】',
    '• 2027Q1-Q2：数值模拟加速与开发战术沙盘智能体全面应用',
    '• 2027Q2-Q3：典型示范区实现油藏-井筒-管网一体化闭环优化',
    '• 2027Q3-Q4：建成"无人值守研究+智能辅助决策"新模式',
    '• 2027Q4：核心业务流程智能化覆盖率>80%',
    '',
    '【风险应对】',
    '• 风险：技术黑箱导致过度依赖，忽略物理规律',
    '• 应对：强制AI输出附带物理解释；保持科研人员批判性思维'
])

# 插入豆包生成的三阶段路线图
add_image_slide(prs, '三阶段实施路线图（2025→2026→2027）', 'chart_roadmap.png')

# 第五章：组织与资源
add_section_divider(prs, '05', '组织保障与预算')

add_content_slide(prs, '组织保障：三级架构', [
    '【领导小组（决策层）】',
    '• 组长：单位分管领导（审批方案、协调资源）',
    '• 副组长：技术/业务部门负责人（审核成果、协调问题）',
    '',
    '【技术组（执行层）】',
    '• 技术负责人：统筹技术方案、把控质量',
    '• 算法工程师（6-8人）：模型微调、智能体开发',
    '• 数据工程师（4-5人）：数据治理、质量核查',
    '',
    '【业务组（需求层）】',
    '• 业务负责人：梳理需求、对接技术、推广应用',
    '• 业务专家（全时4-6人+非全时10人+）：标注数据、验收效果'
])

# 插入豆包生成的组织架构图
add_image_slide(prs, '组织保障架构图（领导小组→技术组+业务组）', 'chart_organization.png')

add_content_slide(prs, '预算框架与算力配置', [
    '【项目总预算：800万元】',
    '• 2025年：320万（硬件120+软件80+人力90+外部30）',
    '• 2026年：300万（硬件80+软件100+人力100+外部20）',
    '• 2027年：180万（硬件50+软件40+人力80+外部10）',
    '',
    '【算力配置】',
    '• 训练集群：利用集团级算力中心（昆仑大模型千卡级GPU）',
    '• 推理集群：本地部署20-30台昇腾910B服务器',
    '• 性能指标：确保推理时延<800ms',
    '',
    '【数据治理工作量】',
    '• 整合3-5套异构数据源',
    '• 解析1万+份历史文档，构建百万级图像样本库',
    '• 构建数万个实体关系的油藏知识图谱'
])

# 第六章：竞争情报
add_section_divider(prs, '06', '竞争情报与差异化')

add_content_slide(prs, '竞争对手能力矩阵', [
    '【华为】',
    '• 核心优势：昇腾算力+盘古底座+"懂行"',
    '• 薄弱环节：行业Know-how深度不足',
    '• 启示：合作+竞争关系，需掌握核心算法能力',
    '',
    '【阿里云】',
    '• 核心优势：云边协同+产学研生态',
    '• 薄弱环节：时序数据处理能力弱于华为',
    '• 启示：可作为补充算力，避免单一供应商绑定',
    '',
    '【百度】',
    '• 核心优势：FM Agent多智能体进化',
    '• 薄弱环节：石油行业渗透浅',
    '• 启示：关注其运筹优化能力'
])

add_content_slide(prs, '我们的差异化优势', [
    '【四大核心优势】',
    '• 数据资产：30万份文献、50TB数据、20万+知识三元组',
    '• 场景深耕：扎根油气开采研究所，懂科研也懂生产',
    '• 自主可控：基于昆仑大模型底座，核心算法自研',
    '• 快速迭代：miniOPC模式，决策链短，响应速度快',
    '',
    '【窗口期判断】',
    '• 国际厂商退出留下2-3年窗口期',
    '• 需快速建立自有能力壁垒',
    '• 与华为、阿里等形成竞合关系'
])

# 第七章：下一步行动
add_section_divider(prs, '07', '下一步行动')

add_content_slide(prs, '行动建议：从即时到中期', [
    '【即时行动（本周）】',
    '• 成立项目筹备组，明确初步人员分工',
    '• 启动数据资产盘点，识别核心数据源',
    '• 与集团数字化部对接，确认昆仑大模型接入方式',
    '',
    '【短期行动（本月）】',
    '• 完成需求调研，明确6大智能体优先级排序',
    '• 制定详细技术方案，确定与华为/阿里云合作边界',
    '• 编制预算细化方案，申请启动资金',
    '',
    '【中期行动（本季度）】',
    '• 完成数据治理攻坚战首批任务',
    '• 搭建开发测试环境',
    '• 启动油藏记忆与类比智能体MVP开发'
])

# 结束页
add_title_slide(prs, 'miniOPC · 一人公司', 'AI赋能 · 轻资产 · 快响应 · 高产出')

# 保存
output_path = "/root/.openclaw/workspace/油气开采研究所AI智能体规划-整合版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
print(f"📄 正文字号：16pt（统一）")
print(f"📝 布局特点：每页最多13个条目，严格分页")
