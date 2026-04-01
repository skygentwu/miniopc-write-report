#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
中石油勘探院油气开采所 - 麦肯锡风格战略分析PPT
基于前期洞察数据生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO

# 视觉规范
PRIMARY_BLUE = RGBColor(0x00, 0x3B, 0x70)
DARK_GRAY = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
BLACK = RGBColor(0x00, 0x00, 0x00)

class OilGasPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        
    def _add_text(self, slide, left, top, width, height, text, 
                  size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font="Arial"):
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.font.name = font
        return shape
        
    def _add_line(self, slide, left, top, width, height, color):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(height))
        line.fill.background()
        line.line.color.rgb = color
        line.line.width = Pt(height)
        return line
    
    def add_cover(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # 主标题
        self._add_text(slide, 1, 2.5, 11, 1.5, 
                      "油气开采所数字化转型战略分析",
                      size=44, bold=True, color=PRIMARY_BLUE, font="Times New Roman")
        
        # 副标题
        self._add_text(slide, 1, 4.2, 11, 0.8,
                      "中石油勘探院油气开采所 · 智能体架构与实施路径",
                      size=24, color=DARK_GRAY)
        
        # 分隔线
        self._add_line(slide, 1, 4.1, 2, 3, PRIMARY_BLUE)
        
        # 日期
        self._add_text(slide, 1, 6.8, 4, 0.3, "2026年3月", size=10, color=LIGHT_GRAY)
        self._add_text(slide, 9, 6.8, 3.333, 0.3, "机密 · 仅供内部讨论使用", 
                      size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
        return slide
    
    def add_section(self, num, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        self._add_text(slide, 1, 2.8, 2, 1.5, f"0{num}" if num < 10 else str(num),
                      size=72, bold=True, color=LIGHT_GRAY, font="Times New Roman")
        self._add_text(slide, 1, 4.3, 11, 1, title,
                      size=28, bold=True, color=PRIMARY_BLUE, font="Times New Roman")
        self._add_line(slide, 1, 5.4, 4, 2, PRIMARY_BLUE)
        return slide
    
    def add_content_slide(self, title, bullets, insights=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        
        # 标题
        self._add_text(slide, 0.5, 0.3, 12, 0.6, title,
                      size=24, bold=True, color=PRIMARY_BLUE, font="Times New Roman")
        self._add_line(slide, 0.5, 0.95, 12.333, 1, LIGHT_GRAY)
        
        # 左侧内容
        content = "\n".join([f"• {b}" for b in bullets])
        self._add_text(slide, 0.5, 1.2, 6, 5.5, content, size=14, color=DARK_GRAY)
        
        # 右侧洞察
        if insights:
            ins_text = "关键洞察\n" + "\n".join([f"→ {i}" for i in insights])
            self._add_text(slide, 7, 1.2, 5.5, 5.5, ins_text, 
                          size=14, color=PRIMARY_BLUE)
        
        # 页码
        self._add_text(slide, 12, 7, 1, 0.3, str(self.slide_count),
                      size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
        return slide
    
    def create_ppt(self):
        """生成完整PPT"""
        # 1. 封面
        self.add_cover()
        
        # 2. 章节1：现状诊断
        self.add_section(1, "油气开采所现状诊断与核心痛点")
        
        self.add_content_slide(
            "业务现状：五大核心业务板块",
            [
                "油藏工程：剩余油分布预测、开发方案优化（关键瓶颈：地质模型精度不足）",
                "采油工程：举升工艺优化、注水方案设计（关键瓶颈：工况诊断依赖经验）",
                "采气工程：气井产能预测、积液诊断（关键瓶颈：多气藏类型差异大）",
                "提高采收率：化学驱方案优化、调剖堵水（关键瓶颈：物理模拟周期长）",
                "数字油气田：数据治理、智能应用（关键瓶颈：系统孤岛严重）"
            ],
            [
                "核心技术需求：机理模型与数据驱动融合",
                "数据资产：20年历史数据，但利用率不足30%",
                "专家依赖：80%决策依赖老师傅经验"
            ]
        )
        
        self.add_content_slide(
            "三大核心痛点分析",
            [
                "痛点一：专家经验传承断层 — 老师傅退休，隐性知识流失，新员工培养周期3-5年",
                "痛点二：数据孤岛与异构 — SCADA、MES、ERP系统互不连通，数据标准不统一",
                "痛点三：决策响应速度慢 — 从问题发现到方案生成平均需2-4周，错失最佳调整窗口"
            ],
            [
                "量化影响：专家依赖导致决策一致性问题",
                "系统割裂导致数据重复录入，效率损失40%",
                "响应延迟导致单井产量损失5-15%"
            ]
        )
        
        # 3. 章节2：智能体架构
        self.add_section(2, "多智能体架构设计方案")
        
        self.add_content_slide(
            "整体架构：云+端双虾协作模式",
            [
                "云端虾（智能大脑）：24小时在线，负责行业情报、案例研究、报告生成",
                "本地虾（执行终端）：部署在院所内网，操作本地系统、访问核心数据、对接现有软件",
                "协作机制：飞书文档作为数据中台，云端生成→本地执行→结果回传",
                "安全设计：数据不出域，敏感操作本地完成，公有云仅处理公开信息"
            ],
            [
                "核心原则：视觉即感知、数据不出域、技能即插件",
                "韧性设计：任一节点故障，任务自动降级或转移"
            ]
        )
        
        self.add_content_slide(
            "六大智能体角色定义",
            [
                "油藏智能体：地质模型自动更新、剩余油预测、开发方案推荐",
                "工况诊断智能体：实时监测、异常预警、故障根因分析",
                "工艺优化智能体：举升参数优化、注水方案生成、能耗分析",
                "知识管理智能体：文献检索、经验沉淀、SOP自动生成",
                "数据治理智能体：数据质量检测、标准对齐、异常清洗",
                "决策支持智能体：多方案对比、风险评估、效果预测"
            ],
            [
                "分工逻辑：各司其职，通过标准化接口协作",
                "渐进部署：从单点智能体开始，逐步扩展协同"
            ]
        )
        
        # 4. 章节3：实施路径
        self.add_section(3, "分阶段实施路径与里程碑")
        
        self.add_content_slide(
            "三阶段实施路线图",
            [
                "第一阶段（1-3月）：数据治理与单点突破 — 统一数据标准，部署工况诊断智能体试点",
                "第二阶段（4-8月）：智能体协同与知识沉淀 — 打通油藏-工艺智能体联动，建立知识库",
                "第三阶段（9-12月）：全局优化与生态构建 — 六大智能体全面协同，输出标准化技能包"
            ],
            [
                "阶段一目标：数据利用率从30%提升至70%",
                "阶段二目标：决策响应时间从2周缩短至3天",
                "阶段三目标：专家经验数字化率超过60%"
            ]
        )
        
        self.add_content_slide(
            "关键成功因素与风险规避",
            [
                "成功因素一：高层支持与组织变革 — 成立智能体推进办公室，明确权责",
                "成功因素二：数据治理先行 — 建立统一数据模型，打通系统孤岛",
                "成功因素三：人机协同设计 — AI推荐+专家确认，避免完全自动化风险",
                "风险规避：数据安全红线、模型可解释性、渐进式替代而非颠覆"
            ],
            [
                "政策契合：符合中石油数字化转型三年行动计划",
                "技术路线：机理+数据双驱动，保证科学性与实用性"
            ]
        )
        
        # 5. 结语页
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_count += 1
        self._add_text(slide, 1, 2.5, 11, 1, 
                      "从数据连接到经验传承，\n从单点优化到全局协同",
                      size=32, bold=True, color=PRIMARY_BLUE, font="Times New Roman")
        self._add_text(slide, 1, 4.5, 11, 1,
                      "油气开采所的智能化转型，\n不是替代专家，而是让专家智慧流动起来",
                      size=20, color=DARK_GRAY)
        
        # 保存
        output = "油气开采所数字化转型战略分析_麦肯锡风格.pptx"
        self.prs.save(output)
        print(f"✅ PPT已生成: {output}")
        print(f"📊 总页数: {self.slide_count}")
        return output

if __name__ == "__main__":
    ppt = OilGasPPT()
    ppt.create_ppt()