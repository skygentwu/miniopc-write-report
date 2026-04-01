#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
顶级咨询公司风格PPT生成器 (McKinsey/BCG/Bain Style)
作者: Kimi Claw
版本: 1.0

使用方式:
    python3 mckinsey_ppt_generator.py --theme "报告主题" --output "输出文件名.pptx"
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np
from io import BytesIO
import argparse
import os

# ============ 视觉规范配置 ============
class VisualStyle:
    """麦肯锡/BCG/贝恩风格视觉规范"""
    
    # 配色方案
    PRIMARY_BLUE = RGBColor(0x00, 0x3B, 0x70)      # 深皇家蓝
    SECONDARY_BLUE = RGBColor(0x00, 0x5A, 0x9E)    # 钴蓝
    DARK_GRAY = RGBColor(0x2C, 0x2C, 0x2C)         # 炭灰
    MEDIUM_GRAY = RGBColor(0x6B, 0x6B, 0x6B)       # 中灰
    LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)        # 浅灰
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    
    # 强调色阶梯
    ACCENT_COLORS = [
        RGBColor(0x00, 0x3B, 0x70),  # 深蓝
        RGBColor(0x4A, 0x90, 0xE2),  # 中蓝
        RGBColor(0x87, 0xCE, 0xEB),  # 浅蓝
        RGBColor(0xB0, 0xB0, 0xB0),  # 灰
        RGBColor(0x6B, 0x6B, 0x6B),  # 深灰
    ]
    
    # 字体配置
    TITLE_FONT = "Times New Roman"  # 高端衬线字体
    BODY_FONT = "Helvetica"          # 极简无衬线
    FALLBACK_FONT = "Arial"
    
    # 字号配置
    TITLE_SIZE = Pt(44)              # 主标题
    SUBTITLE_SIZE = Pt(24)           # 副标题
    HEADING_SIZE = Pt(28)            # 章节标题
    BODY_SIZE = Pt(14)               # 正文
    CAPTION_SIZE = Pt(10)            # 图表标注
    FOOTNOTE_SIZE = Pt(8)            # 脚注
    
    # 布局配置
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)


# ============ PPT生成器核心类 ============
class McKinseyStylePPT:
    """麦肯锡风格PPT生成器"""
    
    def __init__(self, theme_title="战略分析报告"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 宽屏
        self.prs.slide_height = Inches(7.5)
        self.theme_title = theme_title
        self.style = VisualStyle()
        self.slide_count = 0
        
    def _add_text_box(self, slide, left, top, width, height, text, 
                      font_size=Pt(14), font_name="Helvetica", 
                      font_color=None, bold=False, alignment=PP_ALIGN.LEFT):
        """添加文本框"""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        
        if font_color:
            p.font.color.rgb = font_color
        else:
            p.font.color.rgb = self.style.BLACK
            
        return shape
    
    def _add_line(self, slide, left, top, width, height, color, thickness=Pt(1)):
        """添加线条（发丝边框效果）"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        line.fill.background()
        line.line.color.rgb = color
        line.line.width = thickness
        return line
    
    def add_title_slide(self, title=None, subtitle=None, date=None):
        """添加标题页（封面）"""
        if title is None:
            title = self.theme_title
            
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        
        # 主标题 - Times New Roman 高端衬线
        self._add_text_box(
            slide, 
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
            title,
            font_size=self.style.TITLE_SIZE,
            font_name=self.style.TITLE_FONT,
            font_color=self.style.PRIMARY_BLUE,
            bold=True
        )
        
        # 副标题
        if subtitle:
            self._add_text_box(
                slide,
                Inches(1), Inches(4.2), Inches(11.333), Inches(0.8),
                subtitle,
                font_size=self.style.SUBTITLE_SIZE,
                font_name=self.style.BODY_FONT,
                font_color=self.style.MEDIUM_GRAY
            )
        
        # 分隔线
        self._add_line(
            slide,
            Inches(1), Inches(4.1), Inches(2), Pt(2),
            self.style.PRIMARY_BLUE, Pt(3)
        )
        
        # 日期和页脚
        if date:
            self._add_text_box(
                slide,
                Inches(1), Inches(6.8), Inches(4), Inches(0.3),
                date,
                font_size=self.style.CAPTION_SIZE,
                font_name=self.style.BODY_FONT,
                font_color=self.style.LIGHT_GRAY
            )
        
        # 机密标注
        self._add_text_box(
            slide,
            Inches(9), Inches(6.8), Inches(3.333), Inches(0.3),
            "机密 · 仅供内部讨论使用",
            font_size=self.style.CAPTION_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT
        )
        
        return slide
    
    def add_section_divider(self, section_number, section_title):
        """添加章节分隔页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        
        # 章节编号 - 大号数字
        self._add_text_box(
            slide,
            Inches(1), Inches(2.8), Inches(2), Inches(1.5),
            f"0{section_number}" if section_number < 10 else str(section_number),
            font_size=Pt(72),
            font_name=self.style.TITLE_FONT,
            font_color=self.style.LIGHT_GRAY,
            bold=True
        )
        
        # 章节标题
        self._add_text_box(
            slide,
            Inches(1), Inches(4.3), Inches(11.333), Inches(1),
            section_title,
            font_size=self.style.HEADING_SIZE,
            font_name=self.style.TITLE_FONT,
            font_color=self.style.PRIMARY_BLUE,
            bold=True
        )
        
        # 装饰线
        self._add_line(
            slide,
            Inches(1), Inches(5.4), Inches(4), Pt(1),
            self.style.PRIMARY_BLUE, Pt(2)
        )
        
        return slide
    
    def add_content_slide(self, title, bullet_points, insights=None):
        """添加内容页（ bullet points 风格）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        
        # 页面标题
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6),
            title,
            font_size=Pt(24),
            font_name=self.style.TITLE_FONT,
            font_color=self.style.PRIMARY_BLUE,
            bold=True
        )
        
        # 分隔线
        self._add_line(
            slide,
            Inches(0.5), Inches(0.95), Inches(12.333), Pt(1),
            self.style.LIGHT_GRAY, Pt(1)
        )
        
        # 内容区域 - 双栏布局
        content_text = "\n".join([f"• {point}" for point in bullet_points])
        self._add_text_box(
            slide,
            Inches(0.5), Inches(1.2), Inches(6), Inches(5.5),
            content_text,
            font_size=self.style.BODY_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.DARK_GRAY
        )
        
        # 关键洞察框（右侧）
        if insights:
            insight_text = "关键洞察\n" + "\n".join([f"→ {insight}" for insight in insights])
            self._add_text_box(
                slide,
                Inches(7), Inches(1.2), Inches(5.833), Inches(5.5),
                insight_text,
                font_size=self.style.BODY_SIZE,
                font_name=self.style.BODY_FONT,
                font_color=self.style.PRIMARY_BLUE
            )
        
        # 页码
        self._add_text_box(
            slide,
            Inches(12), Inches(7), Inches(1), Inches(0.3),
            str(self.slide_count),
            font_size=self.style.CAPTION_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT
        )
        
        return slide
    
    def add_chart_slide(self, title, chart_type="bar", data=None, 
                       x_label="", y_label="", insights=None):
        """添加图表页 - 使用matplotlib生成高精度图表"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        
        # 设置matplotlib样式
        plt.style.use('seaborn-v0_8-whitegrid')
        plt.rcParams['font.family'] = 'sans-serif'
        plt.rcParams['font.size'] = 10
        plt.rcParams['axes.unicode_minus'] = False
        
        # 配色方案
        colors = ['#003B70', '#4A90E2', '#87CEEB', '#B0B0B0', '#6B6B6B']
        
        fig, ax = plt.subplots(figsize=(10, 5.5), dpi=150)
        fig.patch.set_facecolor('white')
        ax.set_facecolor('white')
        
        if data is None:
            # 示例数据
            data = {
                'categories': ['A', 'B', 'C', 'D', 'E'],
                'values': [85, 72, 90, 68, 78]
            }
        
        if chart_type == "bar":
            bars = ax.bar(data['categories'], data['values'], 
                         color=colors[0], edgecolor='white', linewidth=0.5)
            # 添加数值标签
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height}%', ha='center', va='bottom', fontsize=9)
                       
        elif chart_type == "line":
            ax.plot(data['categories'], data['values'], 
                   color=colors[0], linewidth=2.5, marker='o', markersize=6)
                   
        elif chart_type == "radar":
            # 雷达图
            angles = np.linspace(0, 2*np.pi, len(data['categories']), endpoint=False)
            values = data['values'] + [data['values'][0]]  # 闭合
            angles = np.concatenate((angles, [angles[0]]))
            
            ax = fig.add_subplot(111, projection='polar')
            ax.plot(angles, values, 'o-', linewidth=2, color=colors[0])
            ax.fill(angles, values, alpha=0.25, color=colors[0])
            ax.set_xticks(angles[:-1])
            ax.set_xticklabels(data['categories'])
            
        # 样式调整
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#B0B0B0')
        ax.spines['bottom'].set_color('#B0B0B0')
        ax.tick_params(colors='#6B6B6B')
        ax.set_xlabel(x_label, fontsize=11, color='#2C2C2C')
        ax.set_ylabel(y_label, fontsize=11, color='#2C2C2C')
        
        plt.title(title, fontsize=14, fontweight='bold', color='#003B70', pad=20)
        plt.tight_layout()
        
        # 保存到内存
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight', 
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        plt.close()
        
        # 添加标题
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6),
            title,
            font_size=Pt(24),
            font_name=self.style.TITLE_FONT,
            font_color=self.style.PRIMARY_BLUE,
            bold=True
        )
        
        # 添加图表
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1), 
                                width=Inches(9))
        
        # 添加洞察（右侧）
        if insights:
            insight_text = "关键洞察\n" + "\n".join([f"→ {i}" for i in insights])
            self._add_text_box(
                slide,
                Inches(10), Inches(1.5), Inches(2.833), Inches(5),
                insight_text,
                font_size=self.style.BODY_SIZE,
                font_name=self.style.BODY_FONT,
                font_color=self.style.PRIMARY_BLUE
            )
        
        # 数据来源标注
        self._add_text_box(
            slide,
            Inches(0.5), Inches(6.8), Inches(6), Inches(0.3),
            "数据来源: 内部数据库 · 更新时间: 2026-03-10 · 置信区间: 95%",
            font_size=self.style.FOOTNOTE_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.LIGHT_GRAY
        )
        
        # 页码
        self._add_text_box(
            slide,
            Inches(12), Inches(7), Inches(1), Inches(0.3),
            str(self.slide_count),
            font_size=self.style.CAPTION_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT
        )
        
        return slide
    
    def add_matrix_slide(self, title, matrix_data=None):
        """添加2x2战略矩阵图（竞争定位、BCG矩阵等）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        
        # 标题
        self._add_text_box(
            slide,
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6),
            title,
            font_size=Pt(24),
            font_name=self.style.TITLE_FONT,
            font_color=self.style.PRIMARY_BLUE,
            bold=True
        )
        
        # 使用matplotlib绘制2x2矩阵
        fig, ax = plt.subplots(figsize=(9, 6), dpi=150)
        fig.patch.set_facecolor('white')
        ax.set_facecolor('white')
        
        # 绘制象限分隔线
        ax.axhline(y=0.5, color='#B0B0B0', linewidth=1, linestyle='--')
        ax.axvline(x=0.5, color='#B0B0B0', linewidth=1, linestyle='--')
        
        # 象限标签
        if matrix_data is None:
            matrix_data = {
                'quadrants': ['明星业务\n(高增长/高份额)', '问题业务\n(高增长/低份额)',
                            '现金牛\n(低增长/高份额)', '瘦狗业务\n(低增长/低份额)'],
                'points': [
                    {'x': 0.75, 'y': 0.75, 'label': '业务A', 'size': 300},
                    {'x': 0.25, 'y': 0.75, 'label': '业务B', 'size': 200},
                    {'x': 0.75, 'y': 0.25, 'label': '业务C', 'size': 400},
                    {'x': 0.25, 'y': 0.25, 'label': '业务D', 'size': 150},
                ]
            }
        
        # 绘制气泡
        colors = ['#003B70', '#4A90E2', '#87CEEB', '#B0B0B0']
        for i, point in enumerate(matrix_data['points']):
            ax.scatter(point['x'], point['y'], s=point['size'], 
                      c=colors[i % len(colors)], alpha=0.7, edgecolors='white', linewidth=2)
            ax.annotate(point['label'], (point['x'], point['y']), 
                       xytext=(5, 5), textcoords='offset points',
                       fontsize=10, fontweight='bold', color='#2C2C2C')
        
        # 设置坐标轴
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_xlabel('市场份额 →', fontsize=11, color='#2C2C2C')
        ax.set_ylabel('市场增长率 →', fontsize=11, color='#2C2C2C')
        
        # 移除边框
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#B0B0B0')
        ax.spines['bottom'].set_color('#B0B0B0')
        
        plt.title(title, fontsize=14, fontweight='bold', color='#003B70', pad=20)
        plt.tight_layout()
        
        # 保存
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight',
                   facecolor='white')
        img_stream.seek(0)
        plt.close()
        
        # 添加到幻灯片
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2),
                                width=Inches(8))
        
        # 图例说明
        legend_text = "\n".join(matrix_data['quadrants'])
        self._add_text_box(
            slide,
            Inches(9), Inches(1.5), Inches(3.833), Inches(5),
            legend_text,
            font_size=self.style.BODY_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.DARK_GRAY
        )
        
        # 页码
        self._add_text_box(
            slide,
            Inches(12), Inches(7), Inches(1), Inches(0.3),
            str(self.slide_count),
            font_size=self.style.CAPTION_SIZE,
            font_name=self.style.BODY_FONT,
            font_color=self.style.LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT
        )
        
        return slide
    
    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"✅ PPT已生成: {filename}")
        print(f"📊 总页数: {self.slide_count}")
        return filename


# ============ 使用示例 ============
def create_demo_ppt():
    """创建示例PPT"""
    ppt = McKinseyStylePPT(theme_title="制造业AI转型战略分析")
    
    # 1. 封面
    ppt.add_title_slide(
        title="制造业AI转型战略分析",
        subtitle="从ERP流程分工到智能体协同的演进路径",
        date="2026年3月"
    )
    
    # 2. 章节分隔
    ppt.add_section_divider(1, "市场现状与趋势洞察")
    
    # 3. 内容页
    ppt.add_content_slide(
        title="制造业AI应用现状：三大核心矛盾",
        bullet_points=[
            "AI普及率不足11%，远低于欧美水平，但近八成企业认可AI价值",
            "43%企业尚未部署工业智能体，仅8%实现多场景应用",
            "实验室精度90%+，但产线粉尘、振动下误判率飙升",
            "数据质量不足、数据碎片化、设备兼容性差仍是主要障碍",
            "从单点优化走向全局协同，从感知走向认知是核心趋势"
        ],
        insights=[
            "政策驱动：工信部2027年目标1000个工业智能体",
            "技术拐点：视觉理解+自主执行突破API割裂难题",
            "实施路径：先单点后全局，先质检后排产再维护"
        ]
    )
    
    # 4. 图表页 - 柱状图
    chart_data = {
        'categories': ['AI质检', '预测维护', '智能排产', '数字孪生', '知识图谱'],
        'values': [85, 72, 58, 45, 38]
    }
    ppt.add_chart_slide(
        title="制造业AI应用场景成熟度分布",
        chart_type="bar",
        data=chart_data,
        x_label="应用场景",
        y_label="企业采纳率 (%)",
        insights=[
            "AI质检成熟度最高，已规模化应用",
            "预测维护增速最快，年增长率35%",
            "知识图谱尚处早期，潜力巨大"
        ]
    )
    
    # 5. 2x2矩阵
    matrix_data = {
        'quadrants': [
            '明星：AI质检/视觉检测\n高增长+高渗透，优先投入',
            '问题：预测性维护\n高增长+低渗透，重点培育', 
            '现金牛：MES/ERP集成\n低增长+高渗透，稳健运营',
            '瘦狗：传统规则引擎\n低增长+低渗透，逐步替换'
        ],
        'points': [
            {'x': 0.8, 'y': 0.8, 'label': 'AI质检', 'size': 400},
            {'x': 0.3, 'y': 0.75, 'label': '预测维护', 'size': 300},
            {'x': 0.75, 'y': 0.3, 'label': 'MES集成', 'size': 350},
            {'x': 0.25, 'y': 0.25, 'label': '规则引擎', 'size': 150},
        ]
    }
    ppt.add_matrix_slide(
        title="制造业AI应用战略定位矩阵",
        matrix_data=matrix_data
    )
    
    # 6. 章节分隔
    ppt.add_section_divider(2, "智能体架构与实施路径")
    
    # 7. 内容页 - 架构说明
    ppt.add_content_slide(
        title="多智能体协作架构：六大核心组件",
        bullet_points=[
            "智能大脑/AI调度中心：Master Agent统一指挥，任务拆解与资源分配",
            "数字审计员：安全校验与权限管控，高危操作双智能体验证",
            "人才数字孪生：自动抓取MES操作记录，生成技能画像与智能排班",
            "供应链控制塔：穿透供应商Web门户，实时抓取物流状态",
            "知识工程中心：录屏生成SOP，工艺知识沉淀与经验传承",
            "设备数字孪生：视觉识别老旧机床表盘，预测性维护与异常预警"
        ],
        insights=[
            "核心原则：视觉即感知、数据不出域、技能即插件",
            "部署模式：云端大脑+本地执行，兼顾算力与数据安全",
            "韧性设计：单点故障自动降级，确保产线连续性"
        ]
    )
    
    # 保存
    output_file = "制造业AI转型战略分析_麦肯锡风格.pptx"
    ppt.save(output_file)
    return output_file


# ============ 命令行入口 ============
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="生成麦肯锡风格战略PPT")
    parser.add_argument("--theme", default="战略分析报告", help="报告主题")
    parser.add_argument("--output", default="output.pptx", help="输出文件名")
    parser.add_argument("--demo", action="store_true", help="生成示例PPT")
    
    args = parser.parse_args()
    
    if args.demo:
        create_demo_ppt()
    else:
        # 这里可以扩展为从配置文件或数据文件生成
        ppt = McKinseyStylePPT(theme_title=args.theme)
        ppt.add_title_slide(title=args.theme, date="2026年3月")
        ppt.save(args.output)