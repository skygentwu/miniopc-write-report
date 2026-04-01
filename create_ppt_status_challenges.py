from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿 - 16:9比例
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 浪潮品牌色
INSPUR_BLUE = RGBColor(0, 102, 204)  # 浪潮蓝
INSPUR_ORANGE = RGBColor(255, 102, 0)  # 浪潮橙
INSPUR_RED = RGBColor(255, 51, 0)  # 装饰三角形用
TEXT_DARK = RGBColor(51, 51, 51)  # 正文深灰
TEXT_GRAY = RGBColor(102, 102, 102)  # 副标题灰

# 字体设置
FONT_NAME = '微软雅黑'

def add_decorative_shapes(slide):
    """添加浪潮装饰元素：橙红色大三角形和蓝色小三角形"""
    # 橙红色大三角形（右下角）
    triangle1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2)
    )
    triangle1.fill.solid()
    triangle1.fill.fore_color.rgb = INSPUR_ORANGE
    triangle1.line.fill.background()
    triangle1.rotation = 45
    
    # 蓝色小三角形（右侧）
    triangle2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12), Inches(5.2), Inches(0.8), Inches(0.8)
    )
    triangle2.fill.solid()
    triangle2.fill.fore_color.rgb = INSPUR_BLUE
    triangle2.line.fill.background()
    triangle2.rotation = 30

def add_inspur_logo(slide):
    """添加浪潮Logo区域（左上角）"""
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'inspur 浪潮'
    p.font.name = FONT_NAME
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE

def add_footer_slogan(slide):
    """添加底部标语"""
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
    tf = slogan_box.text_frame
    p = tf.paragraphs[0]
    p.text = '未来，因潮澎湃  Inspur in Future'
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_GRAY

def add_title_slide(prs, title, subtitle=""):
    """封面页 - 标题42号加粗，副标题28号"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加Logo
    add_inspur_logo(slide)
    
    # 主标题 - 42号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题 - 28号
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_NAME
        p2.font.size = Pt(28)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(20)
    
    # 装饰元素
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    
    return slide

def add_section_divider(prs, section_num, section_title):
    """章节分隔页 - 大号章节编号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_inspur_logo(slide)
    
    # 章节编号 - 72号加粗，橙色
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = FONT_NAME
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = INSPUR_ORANGE
    
    # 章节标题 - 36号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    add_decorative_shapes(slide)
    add_footer_slogan(slide)
    return slide

def add_content_slide(prs, title, content_lines):
    """内容页 - 标题32号加粗，二级标题20号加粗，正文18号"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Logo
    add_inspur_logo(slide)
    
    # 页面标题 - 32号加粗
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INSPUR_BLUE
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 判断是二级标题还是正文
        if line.startswith('【') or line.startswith('▎') or line.startswith('■'):
            # 二级标题 - 20号加粗
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = INSPUR_ORANGE
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        elif line.startswith('•') or line.startswith('-') or line.startswith('◆'):
            # 正文列表 - 18号
            p.text = '  ' + line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
        elif line == '':
            # 空行
            p.text = ''
            p.space_after = Pt(8)
        else:
            # 普通正文 - 18号
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
    
    add_footer_slogan(slide)
    return slide

# ==================== PPT内容开始 ====================

# 封面
add_title_slide(prs, '我们的现状与核心挑战', '中石油勘探院油气开采研究所数字化转型调研')

# 第一章：科研侧挑战
add_section_divider(prs, '01', '科研侧挑战：知识断层与效率瓶颈')

add_content_slide(prs, '专家经验传承断层——"人走茶凉"的隐忧', [
    '【核心矛盾】',
    '• 60-70年代参加工作的老专家集中退休，数十年现场经验与"手感"难以显性化留存',
    '',
    '【具体表现】',
    '• 地质解释、油气层识别等高度依赖专家直觉的判断缺乏标准化记录',
    '• 年轻工程师面临陡峭的学习曲线，培养周期长达5-8年',
    '• 同类问题反复摸索，重复"交学费"现象普遍',
    '',
    '【业务影响】',
    '• 关键岗位人才断档风险加剧',
    '• 重大技术决策质量波动',
    '• 院内知识资产流失，难以形成持续积累',
    '',
    '【数据佐证】',
    '• 据行业调研，约40%的石油企业存在核心技术骨干断层风险'
])

add_content_slide(prs, '文献情报处理低效——"淹没在数据里"', [
    '【现状描述】',
    '• 年均新增勘探开发相关论文数万篇、专利数千件',
    '• 行业标准规范持续更新，国际前沿技术快速迭代',
    '',
    '【痛点分析】',
    '• 研究人员日均花费2-3小时在文献检索与阅读上',
    '• 跨语种、跨专业的信息整合困难，知识孤岛现象严重',
    '• 历史研究成果检索困难，"不知道院里有人做过类似工作"',
    '',
    '【深层问题】',
    '• 缺乏统一的领域知识图谱',
    '• 专家经验无法结构化沉淀',
    '• 文献与生产数据割裂，难以形成闭环',
    '',
    '【数据佐证】',
    '• 据内部调研，约35%的科研人员认为"信息过载"是制约创新效率的首要因素'
])

# 第二章：开采侧挑战
add_section_divider(prs, '02', '开采侧挑战：地下黑盒与决策滞后')

add_content_slide(prs, '深层地质条件复杂——"看得见摸不着"', [
    '【技术背景】',
    '• 深层超深层（>6000m）、深水超深水、非常规油气成为勘探开发主战场',
    '• 中石油勘探院正面临"一老、两深、一非"的技术攻关压力',
    '',
    '【核心难题】',
    '• 地下构造复杂，地震资料解释多解性强，储层预测符合率仅65-75%',
    '• 钻井过程中地层不确定性高，井壁失稳、井漏、井喷等风险频发',
    '• 压裂效果难以预测，单井成本动辄数千万甚至上亿元',
    '',
    '【典型场景】',
    '• 万米级超深井：井底温度>200°C，压力>200MPa',
    '• 致密气/页岩气：压裂裂缝网络复杂，返排规律难以预测',
    '• 老油田：储层物性变化大，剩余油分布认识不清',
    '',
    '【行业对标】',
    '• 与国际油服公司相比，复杂井况下的钻井成功率仍有5-10个百分点差距'
])

add_content_slide(prs, '传统数模计算周期长——"算得准等不起"', [
    '【现状分析】',
    '• 油藏数值模拟是开发方案设计的核心工具，但传统工作流存在明显短板',
    '',
    '【瓶颈环节】',
    '• 地质建模到数值模拟的衔接依赖人工干预，数据准备耗时数周',
    '• 单次全油藏精细模拟计算周期长达数天甚至数周',
    '• "方案设计-模拟验证-优化调整"迭代缓慢，难以支撑敏捷决策需求',
    '',
    '【业务后果】',
    '• 错过最佳投产时机',
    '• 在信息不充分情况下"带病"决策',
    '• 增加开发风险与成本',
    '',
    '【典型案例】',
    '• 大庆油田：老井措施方案优化需跨科室协作，审批周期长达2-3周',
    '• 苏里格气田：致密气井生产制度调整滞后于地层压力变化'
])

# 第三章：挑战总结与转型方向
add_section_divider(prs, '03', '挑战总结与数字化转型方向')

add_content_slide(prs, '现状挑战汇总', [
    '【科研侧：知识断层与效率瓶颈】',
    '■ 专家经验传承断层',
    '  • 老专家退休，经验难以显性化留存',
    '  • 年轻工程师培养周期长，重复交学费',
    '',
    '■ 文献情报处理低效',
    '  • 日均2-3小时文献检索',
    '  • 跨专业信息整合困难，知识孤岛严重',
    '',
    '【开采侧：地下黑盒与决策滞后】',
    '■ 深层地质条件复杂',
    '  • 储层预测符合率仅65-75%',
    '  • 钻井风险高，单井成本数千万至上亿元',
    '',
    '■ 传统数模计算周期长',
    '  • 数据准备耗时数周',
    '  • 精细模拟周期长达数天至数周',
    '  • 方案迭代慢，难以支撑敏捷决策'
])

add_content_slide(prs, '数字化转型方向', [
    '【总体思路】',
    '• 构建"数据-算法-业务"三元耦合的认知框架',
    '• 从经验驱动向数据驱动转变',
    '',
    '【科研侧转型路径】',
    '▎领域知识图谱构建',
    '• 结构化沉淀专家经验与历史文献',
    '• 实现跨专业知识的智能检索与关联推荐',
    '',
    '▎智能文献助手',
    '• AI驱动的文献自动摘要与知识抽取',
    '• 多语言文献自动翻译与信息整合',
    '',
    '【开采侧转型路径】',
    '▎实时数模与AI加速',
    '• 基于机器学习的快速代理模型替代传统数值模拟',
    '• 实现分钟级的方案评估与优化',
    '',
    '▎数字孪生与智能预警',
    '• 构建油气藏-井筒-地面全流程数字孪生',
    '• 钻井风险、井筒故障的提前预警与智能处置'
])

add_content_slide(prs, '预期成效', [
    '【科研侧预期成效】',
    '• 专家经验传承：知识图谱覆盖80%以上核心业务领域',
    '• 文献检索效率：AI辅助下检索时间降低70%',
    '• 新人培养周期：从5-8年缩短至3-5年',
    '',
    '【开采侧预期成效】',
    '• 储层预测符合率：从65-75%提升至80-85%',
    '• 数模计算周期：从数周缩短至数小时',
    '• 钻井事故率：降低30-50%',
    '• 单井开发成本：降低10-15%',
    '',
    '【综合效益】',
    '• 研发效率提升30-40%',
    '• 决策响应速度提升5-10倍',
    '• 为中石油勘探院打造"数字中国石油"的标杆示范'
])

# 结束页
add_title_slide(prs, '谢谢', '携手共建智能化油气开采新篇章')

# 保存
output_path = "/root/.openclaw/workspace/中石油勘探院-现状与核心挑战-浪潮模板版.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
print(f"🎨 模板样式：浪潮集团标准模板")
